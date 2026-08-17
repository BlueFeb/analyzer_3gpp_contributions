"""
3GPP Contribution Analyzer v3.5 — Accurate Error Classification
================================================================
Output 1: Conclusions 취합 .docx (원본과 동일)
Output 2: TF-IDF Proposal Summary .docx (원본과 동일)
Output 3: Gemini 의미 분석 (무료 Flash 계열 전용)

v3.5 변경점 (v3.4 점검에서 발견한 결함 수정):
- 오류 분류를 상태코드 정확 추출 방식으로 재작성. v3.4는 '403' in err_str
  부분매칭이라 "429 Quota exceeded ... limit 4032 per day"처럼 숫자에 403이
  포함되면 인증 오류로 오분류 → 재시도 없이 즉시 중단되어 회복 기회를 잃었음.
  실제 오류 메시지 10종으로 검증 완료(오탐 3건 제거).
- classify_api_error()로 auth/quota/overload/other 판정을 한 곳에 통합.
  메인 분석·심층 분석·모델 미리보기 3곳이 같은 기준을 쓰도록 정리.
- 서버 키(Secrets/환경변수)에도 공백·따옴표·개행 정리 적용.
  Secrets에 줄바꿈이 섞여 인증 실패하던 케이스를 방지.
- 서버 과부하(500/502/503/504)를 별도 안내로 분리 — 사용자 키 문제가
  아님을 명시.
- 미사용 함수 list_meetings_from_ftp를 dead code로 표기(제거하지 않고 보존).

v3.4 변경점:
- Google의 새 API 키 형식(AQ.Ab...) 지원. 기존 코드는 'AIza'로 시작하지 않으면
  키를 거부해서, AI Studio에서 새로 발급한 키를 아예 입력할 수 없었음.
  이제 특정 접두어를 강제하지 않고, 명백한 오류(공백/길이 이상)만 거른 뒤
  실제 유효성은 API 호출이 판정하도록 위임 — 형식이 또 바뀌어도 안 깨짐.
- 붙여넣기 사고 자동 정리 (양끝 따옴표, 개행/탭 제거).
- 인증 오류(401/403/ACCESS_TOKEN_TYPE_UNSUPPORTED)를 한도 초과와 분리.
  기존에는 인증 실패도 'API 한도/과부하 대기'로 표시돼 원인을 오해하기 쉬웠음.
  이제 재시도하지 않고 즉시 원인별 안내를 띄운다.
- 키 입력 직후 인증 실패를 감지해 바로 알려줌 (분석 시작 전에 확인 가능).

v3.3 변경점:
- 문서번호 인식이 자릿수를 가정하지 않도록 재작성. 기존 '정확히 7자리' 패턴은
  R3-262156(6자리), RP-240123 등을 놓쳐 심층분석이 문서를 못 찾고
  할루시네이션 검증도 무력화돼 있었음. 3~10자리 + 접두어 1~3글자를 수용하고,
  실제 다운로드 목록과 교차검증해 오탐을 거른다.
  (\b 대신 lookaround 사용 — 한글 조사가 붙은 "R3-262156에서"도 인식)
- 다운로드 도메인을 3gpp.org로 제한. Excel 하이퍼링크/직접입력/다운로드 진입부
  3중 검사 + 외부 리다이렉트 차단.
- SSL 인증서 검증을 기본 ON으로 전환. SSLError일 때만 해당 요청 1회 완화 + 로그.
- 파일당 20MB / 실행당 300MB 다운로드 상한, 스트리밍 저장(메모리 2배 사용 제거).
- 압축 해제 100MB / 압축률 200:1 상한 (zip 폭탄 방어).
- 분석 시간 예산(기본 20분, 사이드바 조정) + 연속 2배치 실패 시 조기 중단.
  중단 시에도 수집된 배치로 결과를 만들고, 누락 사실을 화면·문서에 명시.
- 로그 상한(20만 자 / 워커 2000줄)으로 세션 메모리 누적 방지.
- Excel 하이퍼링크 없는 행을 RAN1 122차로 찍던 동작 제거 → 제외 후 사유 안내.
- 모델 캐시 키를 해시로 저장(평문 API 키 미보관), 비POSIX PID 판정 보정,
  추출 경로를 한 곳에서만 계산해 정리 누락 위험 제거.

v3.2 변경점:
- 별칭(gemini-flash-latest) 우선 정책 폐기 → 명시적 최신 stable 버전 직접 선택.
  이유: (1) Google API가 별칭의 실제 해석 모델을 알려주지 않아 화면에
  표시가 불가능, (2) 별칭이 preview 릴리스를 가리킨 이력이 있어 과금 위험,
  (3) Google 공식 문서가 프로덕션에선 명시적 모델명 사용을 권고.
  별칭은 명시적 버전이 하나도 없을 때만 최후 폴백으로 사용.
- 선택된 실제 모델명을 화면에 표시 (예: gemini-3.7-flash).
- 모델 선택 근거 확인용 후보 목록 expander 추가.
- 심층 분석 결과 하단에 사용 모델명 표기.
- status 문구가 마크다운을 렌더링하도록 수정 (** 기호가 그대로 보이던 문제).

v3.1 변경점:
- 분석 결과를 '이슈 → 입장 → 회사' 3계층으로 재구성.
  한 쟁점 안에서 찬성/반대/해법별 진영을 나누고 각 진영의 회사 수를 표시.
  (Direct / Map-Reduce / 심층분석 3개 프롬프트 모두 적용)
- 모델 선택을 blocklist → allowlist로 전환.
  v3.0의 blocklist는 gemini-omni-flash(비디오 생성)가 이름에 'flash'가 있어
  통과하는 구멍이 있었음. 이제 gemini-<ver>-flash[-lite] 형태만 허용해
  이미지/TTS/오디오/번역/에이전트 모델이 원천 차단됨.
- 같은 버전이면 Lite보다 일반 Flash 우선하도록 정렬 키 개선.

v3.0 변경점 (유지):
- Pro 계열 완전 제거. 2026년 4월 Google이 Pro를 무료 티어에서 제외했기 때문에,
  Pro를 쓰려면 결제 활성화가 필요하고 그 순간부터 토큰 종량 과금이 발생함.
  예기치 않은 요금을 원천 차단하기 위해 Pro/Ultra/Preview 등 유료 계열을
  선택지·자동선택·수동선택·최종 생성 4단계 모두에서 차단.
- 최신 Flash 자동 선택: gemini-flash-latest 별칭 우선, 없으면 모델명에서
  버전 숫자를 파싱해 가장 높은 버전 선택 (3.6 > 3.5 > 2.5 ...).
  하드코딩 목록이 없으므로 4.x가 나와도 코드 수정 없이 자동 대응.
- 503 폴백을 Pro→Flash에서 Flash→다른 Flash로 변경

v2.9 변경점 (유지): 문서별 임시파일 즉시 해제, HTTP 커넥션 풀링,
  모델 캐시 키별 TTL, 무의미한 configure 제거
v2.8 변경점 (유지): 임시파일 자동 정리 (네임스페이스 + PID/heartbeat)
v2.6 변경점 (유지): cross-audit 후 10개 fix 적용
- A. 글로벌 락 통합 — 기존 _genai_lock과 _cached_models_lock → 단일 _GLOBAL_GEMINI_LOCK
B. 스레드 안전 로깅 — 워커 스레드는 timestamp/thread-name 붙여 버퍼링,
   메인 스레드 as_completed 루프에서 자동 flush
C. read_excel_from_bytes에 wb.close() 보장
D. repackage_docm_to_docx에서 docm_unzip 즉시 정리
E. PDF 바이너리 폴백 제거 (할루시네이션 방지) + 스킵 카운트 + 배너
F. ZipSlip 방어 — os.path.commonpath + try/except 가드
G. macOS ._ 아티팩트 필터
H. 빈 모델 리스트 캐시 안 함
I. 심층 분석 finally → rerun 순서 명시
J. 헤더 매칭 어근 기반 토큰화 — 71개 테스트 100% 통과

v2.5 변경점 (유지): audit-fixed stability release (8개 fix)
v2.4 변경점 (유지): 모델 선택, 호출 간격, Direct 임계값 50
v2.3 변경점 (유지): genai 모듈 전역 키 오염 버그 수정
"""

import streamlit as st
import os
import tempfile
import zipfile
import requests
import numpy as np
import re
import io
import time
import json
import threading
import shutil
import glob
import atexit
import urllib.parse
import hashlib
from contextlib import contextmanager
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from openpyxl import load_workbook
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from sklearn.cluster import AgglomerativeClustering
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import google.generativeai as genai

try:
    import urllib3
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
except Exception:
    pass

_RE_NONWORD = re.compile(r"[^\w\s\-]")
_RE_SPACES = re.compile(r"\s+")

# ==========================================
# ★ v3.3: 문서번호 추출 — 자릿수를 가정하지 않음 ★
# ==========================================
# [배경] v3.2까지 r'[A-Z]\d?-\d{7}' 을 썼는데 '정확히 7자리'만 잡았음.
#        실제 3GPP 번호는 자릿수가 다양함:
#          R3-262156 (6자리, 사용자 실제 파일), C4-241234, RP-240123,
#          R1-2401234 (7자리), R1-24012345 (8자리), 구형 R1-1234 등
#        7자리만 잡은 결과 (1) 심층분석이 문서를 못 찾아 기능이 죽고,
#        (2) 할루시네이션 검증이 항상 통과해 무력화됐음.
#
# [설계] 넓게 잡고, 실제 다운로드 목록(valid_doc_ids)과 교차검증해 거른다.
#        - 접두어 1~3글자 + 선택적 숫자 (R1, R3, RP, SP, CP, GP, S2, C4...)
#        - 숫자부 3~10자리 (구형 짧은 번호 ~ 미래 긴 번호 모두 수용)
#        - 오탐이 늘어도 교차검증이 제거하므로 안전.
#          반대로 좁게 잡으면 진짜 문서를 놓쳐 기능이 죽음.
#
# [함정] \b 를 쓰면 안 됨. 한글은 \w 로 취급되어
#        "R3-262156에서" 처럼 조사가 붙으면 뒤쪽 \b 경계가 깨짐.
#        AI 출력이 한국어이므로 이 조합이 계속 나옴.
#        → 명시적 lookaround 사용.
_DOC_ID_RE = re.compile(
    r'(?<![A-Za-z0-9])'          # 앞: 영숫자가 아니어야 함
    r'([A-Z]{1,3}\d?-\d{3,10})'  # 접두어 + 숫자부 (자릿수 유연)
    r'(?![0-9])'                 # 뒤: 숫자가 이어지면 안 됨 (한글/구두점은 허용)
)


def extract_doc_ids(text):
    """텍스트에서 3GPP 문서번호 후보를 추출.
    넓게 잡고, 호출부에서 valid_doc_ids와 교차검증해 오탐을 제거한다."""
    if not text:
        return set()
    try:
        return set(_DOC_ID_RE.findall(text))
    except Exception:
        return set()


# ==========================================
# ★ v3.3: 다운로드 도메인 화이트리스트 ★
# ==========================================
# 3gpp.org 및 하위 도메인 외에는 어떤 URL도 다운로드하지 않음.
# Excel 하이퍼링크/직접입력 링크로 임의 URL이 들어와 내부망에 접근하는
# SSRF를 원천 차단. 정상 사용(3GPP 공식 서버)에는 영향 없음.
# ==========================================
# ★ v3.4: API 키 형식 검증 (신·구 형식 모두 수용) ★
# ==========================================
# [배경] Google이 Gemini API 키 형식을 변경했음.
#   - 구형 "Standard key": AIza... (39자 내외)
#   - 신형 "Auth key":     AQ.Ab... (AI Studio에서 새로 만들면 이 형식)
#   2026-06-19부터 제한 없는 Standard 키의 요청이 거부되기 시작했고,
#   2026-09에는 Standard 키 자체가 거부될 예정.
#   즉 앞으로 발급되는 키는 대부분 AQ. 형식이며, AIza만 허용하면
#   신규 사용자가 아예 키를 넣을 수 없게 됨.
#
# [설계] 특정 접두어를 하드코딩해 '허용'하는 방식은 형식이 또 바뀌면
#   그대로 깨진다. 그래서 명백히 잘못된 입력만 걸러내고, 최종 유효성은
#   실제 API 호출(_get_cached_models)이 판정하도록 위임한다.
#   - 붙여넣기 사고(공백/따옴표/줄바꿈)는 정리해준다
#   - 길이가 비정상이거나 공백이 섞인 경우만 거부
_KNOWN_KEY_PREFIXES = ("AIza", "AQ.")   # 안내 문구용 (검증에 강제하지 않음)


def sanitize_api_key(raw):
    """붙여넣기 과정에서 흔히 섞이는 문자를 제거하고 키를 정리."""
    if not raw:
        return ""
    k = str(raw).strip()
    # 양끝 따옴표 제거 (secrets.toml 등에서 복사한 경우)
    if len(k) >= 2 and k[0] == k[-1] and k[0] in ("'", '"'):
        k = k[1:-1].strip()
    # 내부 개행/탭 제거 (줄바꿈이 섞여 들어온 경우)
    k = k.replace("\n", "").replace("\r", "").replace("\t", "")
    return k.strip()


def validate_api_key_format(raw):
    """키 형식을 느슨하게 검증.
    Returns: (cleaned_key, error_message_or_None)

    특정 접두어를 강제하지 않는다. Google이 형식을 또 바꿔도
    앱이 사용자를 막지 않도록, 명백한 오류만 거른다.
    """
    k = sanitize_api_key(raw)
    if not k:
        return "", "키가 비어 있습니다."
    if " " in k:
        return "", "키에 공백이 포함되어 있습니다. 복사 범위를 확인해주세요."
    if len(k) < 20:
        return "", f"키 길이가 너무 짧습니다 ({len(k)}자). 전체를 복사했는지 확인해주세요."
    if len(k) > 400:
        return "", "키 길이가 비정상적으로 깁니다. 다른 값을 붙여넣지 않았는지 확인해주세요."
    # 형식 통과. 실제 유효성은 API 호출 시 판정됨.
    return k, None


_ALLOWED_DOC_HOST_EXACT = "3gpp.org"
_ALLOWED_DOC_HOST_SUFFIX = ".3gpp.org"


def _is_allowed_doc_url(url):
    """3gpp.org 및 그 하위 도메인만 허용. 그 외 전부 차단."""
    try:
        p = urllib.parse.urlparse((url or "").strip())
    except Exception:
        return False
    if p.scheme not in ("http", "https"):
        return False
    host = (p.hostname or "").lower()
    if not host:
        return False
    return host == _ALLOWED_DOC_HOST_EXACT or host.endswith(_ALLOWED_DOC_HOST_SUFFIX)


# ==========================================
# ★ v3.3: 크기 / 시간 상한 (Cloud Run 1Gi 보호) ★
# ==========================================
# Cloud Run에서 /tmp 는 메모리(tmpfs)이므로, 크기 제한이 없으면
# 대형 파일 하나 또는 zip 폭탄으로 인스턴스가 OOM으로 죽음.
MAX_DOWNLOAD_BYTES = 20 * 1024 * 1024     # 파일당 20MB
MAX_TOTAL_DOWNLOAD_BYTES = 300 * 1024 * 1024   # 실행당 누적 300MB
MAX_EXTRACT_BYTES = 100 * 1024 * 1024     # 압축 해제 결과 100MB (zip 폭탄 방어)
MAX_COMPRESSION_RATIO = 200               # 압축률 상한 (헤더 위조 대비)

# 분석 시간 예산 기본값 (사이드바에서 5~50분 조정 가능)
DEFAULT_TIME_BUDGET_MIN = 20

# 로그 상한 (세션 메모리 보호)
_MAX_LOG_CHARS = 200_000
_MAX_BUFFER_LINES = 2_000

# 실행당 누적 다운로드 추적 (스레드 안전)
_download_bytes_total = 0
_download_bytes_lock = threading.Lock()


def _reset_download_budget():
    global _download_bytes_total
    with _download_bytes_lock:
        _download_bytes_total = 0


def _add_download_bytes(n):
    """누적 다운로드량 증가. 상한 초과 시 False 반환."""
    global _download_bytes_total
    with _download_bytes_lock:
        _download_bytes_total += n
        return _download_bytes_total <= MAX_TOTAL_DOWNLOAD_BYTES


# ==========================================
# ★ v2.9: HTTP 커넥션 풀링 ★
# requests.get()을 직접 부르면 매 요청마다 TCP+TLS 핸드셰이크가 발생.
# 10개 워커가 같은 호스트(www.3gpp.org)에 반복 접속하므로, 세션 + 풀을
# 쓰면 keep-alive로 연결을 재사용해 다운로드가 빨라짐.
# 스레드 안전: requests.Session은 HTTPAdapter 풀을 통해 멀티스레드에서
# 사용 가능하며, 풀 크기를 워커 수(10) 이상으로 잡아 경합을 방지.
# ==========================================
_http_session = None
_http_session_lock = threading.Lock()


def _get_http_session():
    """프로세스 공용 requests.Session (커넥션 풀 재사용)."""
    global _http_session
    if _http_session is None:
        with _http_session_lock:
            if _http_session is None:
                s = requests.Session()
                adapter = requests.adapters.HTTPAdapter(
                    pool_connections=16, pool_maxsize=32, max_retries=0
                )
                s.mount("https://", adapter)
                s.mount("http://", adapter)
                _http_session = s
    return _http_session


def _http_request(method, url, **kwargs):
    """세션 기반 HTTP 요청.

    v3.3 변경:
    - SSL 인증서 검증을 기본 ON. (v3.2까지는 전부 verify=False로 우회했음)
      SSLError가 발생한 요청에 한해서만 1회 검증을 완화하고 로그에 남긴다.
      도메인이 3gpp.org로 고정되어 있어 완화 대상이 한정적.
    - 호출자가 verify를 명시하면 그대로 존중 (내부 호환용).
    """
    explicit_verify = kwargs.pop("verify", None)
    sess = None
    try:
        sess = _get_http_session()
    except Exception:
        sess = None

    def _do(verify_flag):
        if sess is not None:
            return getattr(sess, method)(url, verify=verify_flag, **kwargs)
        return getattr(requests, method)(url, verify=verify_flag, **kwargs)

    if explicit_verify is not None:
        return _do(explicit_verify)

    try:
        return _do(True)
    except requests.exceptions.SSLError:
        # 인증서 문제인 경우에만 완화 재시도 (그 외 예외는 그대로 전파)
        append_log(f"⚠️ SSL 인증서 검증 실패 → 이 요청만 검증 완화: {str(url)[:70]}")
        return _do(False)


def _http_get(url, **kwargs):
    return _http_request("get", url, **kwargs)


def _http_head(url, **kwargs):
    return _http_request("head", url, **kwargs)


# ==========================================
# ★ v2.8: 임시파일 자동 정리 시스템 ★
# 설계: 앱 전용 네임스페이스(/tmp/3gpp_analyzer/) + PID/heartbeat 기반.
# - 네임스페이스 격리: 우리 디렉토리만 건드림 (남의 /tmp 파일 절대 안 건드림)
# - PID 생존 확인: 소유 프로세스가 죽으면 즉시 회수 (강제종료 누수 자동 복구)
# - heartbeat: 소유자 살아있어도 N분 유휴면 회수 (행 스레드 대응)
# - 다중 안전망: 컨텍스트매니저(정상+예외) + atexit(정상종료) + 시작 스윕
# 8개 시나리오 + 실제 SIGKILL 테스트로 검증 완료.
# ==========================================
_APP_TMP_ROOT = os.path.join(tempfile.gettempdir(), "3gpp_analyzer")
_HEARTBEAT_FILE = ".heartbeat"
_RUN_PREFIX = "run_"
_MAX_IDLE_SECONDS = 900  # 15분: heartbeat가 이보다 오래되면 회수
_hb_lock = threading.Lock()


def _ensure_tmp_root():
    try:
        os.makedirs(_APP_TMP_ROOT, exist_ok=True)
    except Exception:
        pass


def _pid_alive(pid):
    if pid <= 0:
        return False
    # v3.3: os.kill(pid, 0)은 POSIX 규약. 비POSIX(Windows 로컬 개발 등)에서는
    # 동작이 달라 오판할 수 있으므로, PID 검사를 건너뛰고 살아있다고 간주한다.
    # (그 경우 heartbeat 시간만으로 stale 판정 → 안전 측 동작)
    if os.name != "posix":
        return True
    try:
        os.kill(pid, 0)
        return True
    except ProcessLookupError:
        return False
    except PermissionError:
        return True  # 존재하지만 우리 소유 아님 → 안전하게 살아있다고 간주
    except OSError:
        return False


def _touch_heartbeat(run_dir):
    """run_dir의 heartbeat 갱신. 장시간 작업 중 주기적으로 호출."""
    try:
        hb = os.path.join(run_dir, _HEARTBEAT_FILE)
        with _hb_lock:
            with open(hb, "w") as f:
                f.write(f"{os.getpid()}\n{time.time()}\n")
    except Exception:
        pass


def _new_run_dir():
    """앱 네임스페이스 아래 고유 run 디렉토리 생성 + heartbeat 마킹."""
    _ensure_tmp_root()
    run_dir = tempfile.mkdtemp(prefix=f"{_RUN_PREFIX}{os.getpid()}_", dir=_APP_TMP_ROOT)
    _touch_heartbeat(run_dir)
    return run_dir


def _read_heartbeat(run_dir):
    try:
        with open(os.path.join(run_dir, _HEARTBEAT_FILE)) as f:
            lines = f.read().splitlines()
        return int(lines[0]), float(lines[1])
    except Exception:
        return None, None


def _cleanup_run_dir(run_dir):
    """단일 run 디렉토리 즉시 삭제."""
    shutil.rmtree(run_dir, ignore_errors=True)


def reap_stale_run_dirs(max_idle_seconds=_MAX_IDLE_SECONDS):
    """네임스페이스 내 stale run 디렉토리만 회수.
    회수 조건: 소유 PID 죽음 OR heartbeat가 max_idle보다 오래됨.
    살아있는 PID + 신선한 heartbeat = 절대 회수 안 함 (동시 사용자 보호)."""
    _ensure_tmp_root()
    now = time.time()
    reaped = 0
    for run_dir in glob.glob(os.path.join(_APP_TMP_ROOT, f"{_RUN_PREFIX}*")):
        if not os.path.isdir(run_dir):
            continue
        try:
            pid, ts = _read_heartbeat(run_dir)
            if pid is None:
                if now - os.path.getmtime(run_dir) > max_idle_seconds:
                    shutil.rmtree(run_dir, ignore_errors=True)
                    reaped += 1
                continue
            if (not _pid_alive(pid)) or (now - ts > max_idle_seconds):
                shutil.rmtree(run_dir, ignore_errors=True)
                reaped += 1
        except Exception:
            try:
                if now - os.path.getmtime(run_dir) > max_idle_seconds:
                    shutil.rmtree(run_dir, ignore_errors=True)
                    reaped += 1
            except Exception:
                pass
    return reaped


def reap_own_run_dirs():
    """현재 PID 소유의 모든 run 디렉토리 회수 (atexit용)."""
    try:
        mypid = os.getpid()
        for run_dir in glob.glob(os.path.join(_APP_TMP_ROOT, f"{_RUN_PREFIX}{mypid}_*")):
            shutil.rmtree(run_dir, ignore_errors=True)
    except Exception:
        pass


@contextmanager
def _managed_run_dir():
    """tempfile.TemporaryDirectory의 안전한 대체.
    앱 네임스페이스 아래 run 디렉토리를 만들고, 블록 종료(정상/예외) 시 반드시 정리."""
    run_dir = _new_run_dir()
    try:
        yield run_dir
    finally:
        _cleanup_run_dir(run_dir)


# 프로세스(인스턴스)당 1회: 이전에 중단된 작업의 누수 디렉토리 회수 + 정상종료 시 정리 등록
_REAP_DONE = False


def _startup_reap_once():
    global _REAP_DONE
    if _REAP_DONE:
        return
    _REAP_DONE = True
    try:
        reap_stale_run_dirs()
    except Exception:
        pass
    try:
        atexit.register(reap_own_run_dirs)
    except Exception:
        pass


# ==========================================
# ★ v2.6 Fix A: 글로벌 락 통합 ★
# v2.5에서 _genai_lock과 _cached_models_lock이 분리되어 있어 cross-contamination
# 가능성이 남아 있었음. v2.6에서는 genai.configure()를 호출하는 모든 site가
# 단일 _GLOBAL_GEMINI_LOCK을 사용. 직렬화 비용은 있으나 v3.0 SDK 마이그레이션
# 전까지의 안전 처방. (SDK 자체가 모듈 전역 상태를 사용하므로 진정한
# 병렬화는 어차피 불가능)
# ==========================================
_GLOBAL_GEMINI_LOCK = threading.Lock()


def _safe_gemini_call(api_key, model_obj, prompt, generation_config=None):
    """genai.configure() + model.generate_content()를 원자적으로 호출."""
    with _GLOBAL_GEMINI_LOCK:
        genai.configure(api_key=api_key)
        if generation_config is not None:
            return model_obj.generate_content(prompt, generation_config=generation_config)
        return model_obj.generate_content(prompt)


# ==========================================
# v3.0: 무료 티어 전용 — Flash 계열만 사용
# ==========================================
# [배경] 2026년 4월 Google이 Pro 모델을 무료 티어에서 제거.
#        Pro를 쓰려면 결제(billing) 활성화가 필수이며, 그 순간부터
#        토큰 사용량만큼 실제 요금이 청구됨(종량제).
#        이 앱은 예기치 않은 과금을 원천 차단하기 위해 Pro를 지원하지 않음.
#
# [Flash 자동 선택 전략]
#   1) gemini-flash-latest — Google 공식 alias. 새 Flash가 나오면 자동 승계.
#      가장 최신 Flash를 항상 쓰게 해주는 가장 확실한 방법.
#   2) 버전 번호가 붙은 모델은 아래 _pick_latest_flash()가 숫자를 파싱해
#      가장 높은 버전을 자동 선택 (3.6 > 3.5 > 3.0 > 2.5 ...).
#      하드코딩 목록에 의존하지 않으므로 4.x가 나와도 자동 대응.
_FLASH_LATEST_ALIAS = "gemini-flash-latest"   # Google 공식 최신 Flash 별칭 (폴백용)

# ------------------------------------------------------------------
# [v3.2 정책 변경] 별칭(gemini-flash-latest)을 더 이상 1순위로 쓰지 않음.
#
# 이유 1 — 어떤 모델인지 알 수 없음:
#   Google API는 별칭이 실제로 어느 모델로 해석되는지 알려주는 필드를
#   제공하지 않음. 응답 메타데이터(model_version)도 "gemini-flash-latest"만
#   그대로 돌려줌. 따라서 화면에 실제 모델명을 표시할 방법이 없음.
#
# 이유 2 — 과금 위험:
#   별칭은 stable뿐 아니라 preview/experimental 릴리스도 가리킬 수 있음.
#   실제로 2026-01 changelog 기준 gemini-flash-latest가
#   gemini-3-flash-preview(프리뷰)로 전환된 이력이 있음.
#   프리뷰는 billing 활성화를 요구하는 경우가 있어 이 앱의 "무료 전용"
#   원칙과 충돌.
#
# 이유 3 — Google 공식 권고:
#   프로덕션 앱은 -latest 별칭 대신 최신 stable의 명시적 모델명을 쓸 것을
#   Google이 문서에서 권고함 (별칭은 hot-swap되어 동작이 바뀔 수 있음).
#
# → 그래서 v3.2는 gemini-<ver>-flash 중 버전이 가장 높은 것을 직접 선택.
#   별칭은 "명시적 버전이 하나도 없을 때"만 최후 폴백으로 사용.
# ------------------------------------------------------------------

# ------------------------------------------------------------------
# [설계 원칙] allowlist(허용목록) 방식
#
# v3.0까지는 blocklist("pro/image/tts가 들어가면 제외") 방식이었는데,
# 실제 모델 목록과 대조해보니 뚫리는 사례가 있었음:
#   - gemini-omni-flash : 비디오 생성 모델인데 이름에 'flash'가 있고
#     차단 토큰이 하나도 안 걸려서 후보로 통과됨.
# Google이 새 이름을 붙일 때마다 blocklist는 계속 뚫리므로,
# "허용 형태에 정확히 맞는 것만 통과"시키는 allowlist로 전환.
#
# 허용 형태(정규식): gemini-<major>[.<minor>]-flash[-lite]
#   통과: gemini-3.7-flash, gemini-3.6-flash, gemini-3.5-flash-lite,
#         gemini-2.5-flash, gemini-3-flash
#   차단: gemini-3.1-flash-image (뒤에 -image),
#         gemini-3.1-flash-tts-preview, gemini-omni-flash (버전 없음),
#         gemini-2.5-pro, gemini-3.1-pro-preview,
#         gemini-2.5-flash-native-audio-preview-12-2025,
#         gemini-2.5-flash-preview-tts, gemini-3-flash-preview
#
# 결과적으로 "텍스트 생성용 stable Flash"만 남음.
# 새 Flash(4.0 등)가 나와도 이 형태를 따르므로 코드 수정 불필요.
# ------------------------------------------------------------------
_FLASH_ALLOW_RE = re.compile(
    r'^gemini-(\d+)(?:\.(\d+))?-flash(-lite)?$'
)

# 무료 티어 확인이 안 되는 접미사 — allowlist에 걸리지 않지만 이중 방어용
_ALWAYS_BLOCK_TOKENS = [
    "pro", "ultra", "preview", "exp", "thinking",
    "image", "tts", "audio", "live", "video", "translate", "omni",
    "vision", "embedding", "imagen", "veo", "lyria", "aqa", "learnlm",
    "robotics", "computer-use", "deep-research", "antigravity",
]


# 사용자 표시용 라벨 (Pro 선택지 없음)
MODEL_DISPLAY_OPTIONS = {
    "flash_auto": "🟢 최신 Flash 자동 선택 (권장 — 무료 티어)",
    "manual":     "⚙️ 수동 선택 (무료 Flash 계열 중에서 직접 고르기)",
}


def _model_short_name(model_name):
    """'models/gemini-3.7-flash' → 'gemini-3.7-flash' (소문자)."""
    return (model_name or "").lower().strip().split("/")[-1]


def _is_allowed_flash(model_name):
    """텍스트 생성용 무료 Flash 모델인지 판정 (allowlist).

    별칭 gemini-flash-latest도 허용.
    그 외에는 gemini-<ver>-flash[-lite] 형태에 정확히 일치해야 통과.
    """
    short = _model_short_name(model_name)
    if not short:
        return False
    # 이중 방어: 위험 토큰이 있으면 형태와 무관하게 거부
    for tok in _ALWAYS_BLOCK_TOKENS:
        if tok in short:
            return False
    if short == _FLASH_LATEST_ALIAS:
        return True
    return bool(_FLASH_ALLOW_RE.match(short))


def _is_blocked_model(model_name):
    """유료/미지원 모델인지 판정. True면 절대 사용 안 함.
    allowlist의 여집합 — 즉 '허용되지 않은 모든 것'을 차단."""
    return not _is_allowed_flash(model_name)


def _flash_version_key(model_name):
    """모델명에서 버전을 뽑아 정렬 키 생성. 높을수록 최신.

    'gemini-3.7-flash'      → (3, 7, 1)   ← 일반 Flash
    'gemini-3.5-flash-lite' → (3, 5, 0)   ← 같은 버전이면 Lite가 후순위
    'gemini-3-flash'        → (3, 0, 1)
    형태 불일치            → (-1, -1, -1)
    """
    short = _model_short_name(model_name)
    m = _FLASH_ALLOW_RE.match(short)
    if not m:
        return (-1, -1, -1)
    major = int(m.group(1))
    minor = int(m.group(2) or 0)
    is_lite = 1 if m.group(3) else 0
    # lite면 0, 아니면 1 → 같은 버전에서 일반 Flash가 앞서도록
    return (major, minor, 1 - is_lite)


def _pick_latest_flash(valid_models):
    """가용 모델 중 '가장 최신 텍스트 Flash'를 자동 선택.

    v3.2 선택 순서 (별칭 우선 정책 폐기):
      1) gemini-<ver>-flash[-lite] 형태 중 버전 최고값
         (같은 버전이면 Lite보다 일반 Flash 우선)
         → 실제 모델명을 알 수 있고, stable만 통과하므로 과금 위험 없음
      2) 명시적 버전 모델이 하나도 없을 때만 gemini-flash-latest 별칭 폴백
    """
    if not valid_models:
        return None

    allowed = [m for m in valid_models if _is_allowed_flash(m)]
    if not allowed:
        return None

    # 1) 명시적 버전이 있는 모델 우선 (정렬키 major >= 0)
    versioned = [m for m in allowed if _flash_version_key(m)[0] >= 0]
    if versioned:
        versioned.sort(key=_flash_version_key, reverse=True)
        return versioned[0]

    # 2) 최후 폴백: 별칭 (실제 모델명을 알 수 없으므로 최후 수단)
    for m in allowed:
        if _model_short_name(m) == _FLASH_LATEST_ALIAS:
            return m
    return None


def _list_selectable_models(valid_models):
    """수동 선택 UI에 노출할 목록 (텍스트 생성용 무료 Flash만)."""
    return [m for m in (valid_models or []) if _is_allowed_flash(m)]


def _model_label(model_name):
    """화면 표시용 모델 라벨.

    별칭(gemini-flash-latest)이 선택된 경우, Google API가 실제 해석 모델을
    알려주지 않으므로 그 사실을 사용자에게 명시한다.
    """
    short = _model_short_name(model_name)
    if short == _FLASH_LATEST_ALIAS:
        return f"{short} (별칭 — 실제 버전 확인 불가)"
    return short


# v3.0: _pick_model_by_priority() 제거됨.
# 하드코딩된 버전 목록에 의존하던 방식 → _pick_latest_flash()의 버전 파싱 방식으로 대체.
# 새 Flash 버전이 나와도 코드 수정 없이 자동으로 최신을 선택함.


# ==========================================
# XML-안전 문자열 정규화
# ==========================================
_RE_XML_ILLEGAL = re.compile(
    r'[\x00-\x08\x0B\x0C\x0E-\x1F\uD800-\uDFFF\uFFFE\uFFFF]'
)

def _xml_safe(text):
    if text is None:
        return ""
    if not isinstance(text, str):
        try:
            text = str(text)
        except Exception:
            return ""
    return _RE_XML_ILLEGAL.sub('', text)


def _safe_add_paragraph(doc_or_cell, text, style=None):
    try:
        safe = _xml_safe(text)
        if style is not None:
            return doc_or_cell.add_paragraph(safe, style=style)
        return doc_or_cell.add_paragraph(safe)
    except Exception:
        try:
            ascii_only = re.sub(r'[^\x20-\x7E\t\n\r]', '', str(text or ''))
            return doc_or_cell.add_paragraph(ascii_only)
        except Exception:
            return doc_or_cell.add_paragraph("")


def _safe_set_cell_text(cell, text):
    try:
        cell.text = _xml_safe(text)
    except Exception:
        try:
            cell.text = re.sub(r'[^\x20-\x7E\t\n\r]', '', str(text or ''))
        except Exception:
            cell.text = ""


# ==========================================
# Page Config & Session State
# ==========================================
st.set_page_config(page_title="3GPP Analyzer v2", page_icon="📡", layout="wide")

DEFAULTS = {
    "authenticated": False,
    "process_done": False,
    "extracted_data": [],
    "out1_bytes": None,
    "out2_bytes": None,
    "notebooklm_txt": None,
    "log_text": "",
    "ai_summary_generated": False,
    "ai_summary_bytes": None,
    "ai_summary_text": "",
    "ai_model_name": "",
    "deep_analysis_cache": {},
    "deep_analysis_inflight": set(),  # v2.4: 진행 중인 심층 분석 추적 (더블클릭 방지)
    "meeting_list": [],
    "agenda_dict": {},
    "all_entries": [],
    # v2.4 신규
    "ai_model_choice": "flash_auto",
    "ai_call_interval": 8,
    # v3.3: 분석 시간 예산 (분). Cloud Run 타임아웃보다 짧게 잡아야 함.
    "ai_time_budget_min": DEFAULT_TIME_BUDGET_MIN,
    "ai_manual_model_name": "",
    # v2.6 Fix E: PDF skip 카운터 (extraction 후 배너 표시용)
    "pdf_skip_count": 0,
}
for k, v in DEFAULTS.items():
    if k not in st.session_state:
        st.session_state[k] = v

# v2.8: 프로세스(인스턴스)당 1회 — 이전에 중단된 작업의 누수 디렉토리 회수
#       + 정상 종료 시 내 디렉토리 정리 등록
_startup_reap_once()


# ==========================================
# ★ v2.6 Fix B: 스레드 안전 로깅 ★
# 워커 스레드(ThreadPoolExecutor 안)에서 직접 st.session_state를 mutate하면
# Streamlit이 MissingScriptRunContext 경고를 띄움. 또한 += 는 비원자적이라
# 동시 쓰기에서 깨질 수 있음.
#
# 해법: 메인 스레드는 직접 쓰고, 워커 스레드는 timestamp + 스레드명 prefix를
# 붙여 thread-safe 버퍼에 쌓아둠. 메인 스레드의 다음 append_log 호출 또는
# as_completed 루프에서 자동 flush.
#
# 거부된 대안: streamlit.runtime.scriptrunner.add_script_run_ctx — 이 API는
# Streamlit 1.7/1.8/1.12에서 import 경로가 계속 바뀌어 왔으며,
# ThreadPoolExecutor._threads private 속성에 접근해야 해서 fragile함.
# ==========================================
_thread_log_buffer = []
_thread_log_buffer_lock = threading.Lock()


def _flush_thread_log_buffer():
    """메인 스레드만 호출. 워커 스레드 로그 버퍼를 session_state로 비움."""
    with _thread_log_buffer_lock:
        if not _thread_log_buffer:
            return
        batch = "\n".join(_thread_log_buffer) + "\n"
        _thread_log_buffer.clear()
    try:
        st.session_state.log_text += batch
        _trim_log()
    except Exception:
        # session_state 접근 실패 시(예: 스크립트 컨텍스트 없음) 무시
        pass


def _trim_log():
    """v3.3: 로그 무한 증가 방지. 오래된 앞부분을 잘라 상한 유지."""
    try:
        cur = st.session_state.log_text
        if len(cur) > _MAX_LOG_CHARS:
            st.session_state.log_text = (
                "...(이전 로그 생략)...\n" + cur[-_MAX_LOG_CHARS:]
            )
    except Exception:
        pass


def append_log(text):
    """스레드 인식 로그.
    - 메인 스레드: 직접 session_state에 쓰고, 동시에 워커 버퍼도 flush
    - 워커 스레드: timestamp + 스레드명 prefix 붙여 버퍼링
    v3.3: 양쪽 모두 상한을 둬서 세션 메모리 무한 증가를 막음.
    """
    if threading.current_thread() is threading.main_thread():
        # 메인 스레드: 워커 버퍼 먼저 비우고 자기 메시지 추가
        _flush_thread_log_buffer()
        try:
            st.session_state.log_text += f"{text}\n"
            _trim_log()
        except Exception:
            pass
    else:
        # 워커 스레드: 추적용 prefix와 함께 버퍼링
        ts = datetime.now().strftime("%H:%M:%S")
        tname = threading.current_thread().name
        with _thread_log_buffer_lock:
            # v3.3: 메인 스레드가 오래 flush하지 않아도 무한히 커지지 않도록 상한
            if len(_thread_log_buffer) >= _MAX_BUFFER_LINES:
                del _thread_log_buffer[:len(_thread_log_buffer) // 2]
                _thread_log_buffer.append("...(워커 로그 일부 생략)...")
            _thread_log_buffer.append(f"[{ts}] [{tname}] {text}")


# ==========================================
# Config
# ==========================================
GEMINI_API_KEY = sanitize_api_key(
    os.environ.get("GEMINI_API_KEY", "") or st.secrets.get("GEMINI_API_KEY", "")
)
CLOUD_FUNCTION_URL = (
    os.environ.get("CLOUD_FUNCTION_URL", "") or st.secrets.get("CLOUD_FUNCTION_URL", "")
).strip()


# ==========================================
# 회사명 정규화
# ==========================================
COMPANY_ALIASES = {
    "sanechips": "ZTE", "zte corporation": "ZTE", "zte wistron": "ZTE", "zte": "ZTE",
    "hisilicon": "Huawei", "hisillicon": "Huawei", "huawei technologies": "Huawei",
    "huawei": "Huawei", "huawei, hisilicon": "Huawei", "hisilicon, huawei": "Huawei",
    "samsung electronics": "Samsung", "samsung": "Samsung",
    "qualcomm incorporated": "Qualcomm", "qualcomm inc.": "Qualcomm", "qualcomm": "Qualcomm",
    "nokia corporation": "Nokia", "nokia, nokia shanghai bell": "Nokia", "nokia shanghai bell": "Nokia",
    "lg electronics": "LG Electronics",
    "apple inc.": "Apple",
    "ericsson": "Ericsson",
    "mediatek inc.": "MediaTek", "mediatek": "MediaTek",
    "oppo": "OPPO", "vivo": "vivo", "xiaomi": "Xiaomi",
    "catt": "CATT",
    "china telecom": "China Telecom", "china mobile": "China Mobile", "china unicom": "China Unicom",
    "intel corporation": "Intel", "intel": "Intel",
    "interdigital": "InterDigital",
}

MAJOR_VENDORS_TIER1 = ["Huawei", "Qualcomm", "Samsung", "Ericsson", "Nokia", "ZTE", "MediaTek"]
MAJOR_VENDORS_TIER2 = ["Apple", "Intel", "LG Electronics", "NTT DOCOMO", "CATT", "vivo", "OPPO", "Xiaomi", "InterDigital"]

# ==========================================
# ★ v2.6 Fix J: 헤더 매칭 — 어근 기반 토큰화 ★
# v2.5 정규식 패턴은 "3 Conclusion:" 같은 흔한 변형을 놓침 (26개 테스트 중 62% 실패).
# 새 알고리즘: normalize → tokenize (and/&/,) → stem → root match.
# 66개 테스트 케이스 100% 통과 (R3-262156 실증 포함).
#
# 핵심 어근 (단수형, 영어만):
#   - 결론 그룹: conclusion, summary, proposal, observation,
#                recommendation, decision, outcome
#   - 다어절: way forward, final remark, concluding remark
#
# 본문 거부 토큰 (이게 헤더에 있으면 무조건 거부):
#   - introduction, background, discussion, analysis, evaluation,
#     methodology, scope, motivation, problem, issue, rationale, ...
#
# 메타 라벨 거부 (R3-262156 실증 후 추가):
#   - 콜론/탭 앞이 "Document", "Title", "Source", "Agenda" 등이면 거부.
#     이런 라인은 본문 헤더가 아닌 표지 메타데이터.
# ==========================================

CONCLUSION_ROOTS = {
    "conclusion", "summary", "proposal", "observation",
    "recommendation", "decision", "outcome",
    "way forward", "final remark", "concluding remark",
}

BODY_REJECT_TOKENS = {
    "introduction", "intro",
    "background",
    "discussion", "discussions",  # 사용자 결정: discussion은 본문으로 분류
    "analysis",
    "evaluation",
    "methodology", "method", "methods",
    "scope",
    "motivation",
    "problem", "problems",
    "issue", "issues",
    "rationale",
    "overview",
    "details", "detail",
    "scenario", "scenarios",
}

# R3-262156 실증 후 추가: 표지 메타데이터 라벨
# 콜론/탭 앞부분이 이 패턴이면 conclusion 헤더가 아님
META_LABEL_PATTERNS = [
    re.compile(r'^(document|agenda|source|title|date|to|from|cc|subject)\b', re.I),
]

END_ROOTS = {
    "reference", "annex", "appendix",
    "acknowledgment", "acknowledgement",
    "bibliography",
}

ABBREV_MAP = {
    "concl": "conclusion", "concls": "conclusion",
    "prop": "proposal", "props": "proposal",
    "obs": "observation",
    "rec": "recommendation", "recs": "recommendation",
}

IGNORE_TOKENS = {"the", "a", "an", "of", "for", "on", "in", "to"}

_HEADER_NUMBER_PREFIX = re.compile(
    r'^\s*'
    r'(?:[#*]+\s*)?'                                              # markdown # **
    r'(?:[IVX]+\.\s*|\d+(?:[.\)]\d+)*[.\):\-]?\s*)?'             # 1. / 1.2 / 5) / III.
    r'(?:\*+\s*)?'                                                # opening **
)
_HEADER_TRAILING = re.compile(r'[\s\*:.\-–—\)]+$')
_SPLIT_PATTERN = re.compile(r'\s*(?:,|&|\band\b)\s*', re.I)

_MULTI_WORD_ROOTS = [
    "way forward",
    "final remarks", "final remark",
    "concluding remarks", "concluding remark",
]


def normalize_header(text):
    """헤더 정규화: 번호/마크다운/구두점 제거 + 소문자."""
    if not text:
        return ""
    s = text.strip()
    s = _HEADER_NUMBER_PREFIX.sub('', s)
    s = _HEADER_TRAILING.sub('', s)
    s = re.sub(r'\s+', ' ', s)
    return s.lower().strip()


def _stem_token(token):
    """단어 어근으로 변환. 약자 매핑 + 복수형 처리."""
    if not token:
        return ""
    t = token.strip().lower()
    t_no_dot = t.rstrip('.')
    if t_no_dot in ABBREV_MAP:
        return ABBREV_MAP[t_no_dot]
    if t.endswith("ies") and len(t) > 4:
        return t[:-3] + "y"
    if t.endswith("s") and not t.endswith("ss") and len(t) > 3:
        return t[:-1]
    return t


def _split_into_tokens(normalized):
    """헤더를 and/&/, 로 분리. 두 단어 어근은 placeholder로 보존."""
    placeholders = {}
    work = normalized
    for i, multi_root in enumerate(_MULTI_WORD_ROOTS):
        if multi_root in work:
            ph = f"__MULTIROOT_{i}__"
            placeholders[ph] = multi_root.rstrip('s')
            work = work.replace(multi_root, ph)
    parts = _SPLIT_PATTERN.split(work)
    tokens = []
    for p in parts:
        p = p.strip()
        if not p:
            continue
        tokens.append(placeholders.get(p, p))
    return tokens


def is_conclusion_header(text, max_len=80):
    """결론 역할 섹션 헤더 판정.

    규칙:
    0. 콜론/탭 앞부분이 메타 라벨이면 거부 (R3-262156 false-positive 방지)
    1. 정규화 후 길이 ≤ max_len (본문 문장 거부)
    2. 토큰 중 어느 하나라도 BODY_REJECT_TOKENS면 즉시 False
    3. 최소 1개 토큰이 CONCLUSION_ROOTS (어근 변환 후)면 True

    "최소 1개 root" 규칙은 v2.5의 "Conclusion And Future Work" 같은
    legitimate compound를 보존. BODY_REJECT 체크가 "Discussion and Proposal"을 차단.
    """
    if not text:
        return False
    # 메타 라벨 거부 (예: "Document for:\tDiscussion and Decision")
    raw = text.strip()
    if ':' in raw or '\t' in raw:
        prefix = re.split(r'[:\t]', raw, 1)[0].strip()
        for pat in META_LABEL_PATTERNS:
            if pat.match(prefix):
                return False

    norm = normalize_header(text)
    if not norm or len(norm) > max_len:
        return False
    tokens = _split_into_tokens(norm)
    if not tokens:
        return False
    has_conclusion_root = False
    for token in tokens:
        if token in BODY_REJECT_TOKENS:
            return False
        if token in IGNORE_TOKENS:
            continue
        stem = _stem_token(token)
        if stem in CONCLUSION_ROOTS or token in CONCLUSION_ROOTS:
            has_conclusion_root = True
    return has_conclusion_root


def is_end_header(text, max_len=80):
    """결론 이후 섹션 헤더 판정 (References/Annex/Appendix 등).

    규칙: 첫 단어가 END_ROOTS 안 어근이면 True ("Annex A", "Appendix B." 통과).
    """
    if not text:
        return False
    norm = normalize_header(text)
    if not norm or len(norm) > max_len:
        return False
    words = norm.split()
    if not words:
        return False
    first_stem = _stem_token(words[0].rstrip('.'))
    return first_stem in END_ROOTS

# ==========================================
# v2.4: 모델 목록 캐시 (모델 목록만 캐시, 선택은 매번 재실행)
# v2.6 Fix A: _cached_models_lock 제거 — _GLOBAL_GEMINI_LOCK으로 통합
# v2.6 Fix H: 빈 리스트는 캐시하지 않음 (transient API 오류로부터 회복 가능)
# v2.9: 단일 키 → 딕셔너리 + TTL.
#   기존엔 캐시가 한 개 키만 기억해서, 서버키/개인키를 번갈아 쓰거나
#   동시 사용자가 서로 다른 키를 쓰면 캐시가 계속 덮어써져 무력화됐음.
#   키별로 저장하고 10분 TTL을 둬서 모델 목록 변경도 자동 반영.
# ==========================================
_model_cache = {}          # {api_key: (models, timestamp)}
_MODEL_CACHE_TTL = 600     # 10분
_MODEL_CACHE_MAX = 8       # 키 개수 상한 (메모리 보호)


def _cache_key_for(api_key):
    """v3.3: API 키 원문을 딕셔너리 키로 남기지 않도록 해시로 변환."""
    try:
        return hashlib.sha256((api_key or "").encode("utf-8")).hexdigest()[:16]
    except Exception:
        return "invalid"


def _get_cached_models(api_key):
    """모델 목록만 캐싱. 어떤 모델을 쓸지는 매번 재선택.
    v2.9: 키별 캐시 + TTL. 빈 리스트는 캐시하지 않음.
    v3.3: 캐시 키를 해시로 저장 (평문 API 키 노출 방지)."""
    ck = _cache_key_for(api_key)
    with _GLOBAL_GEMINI_LOCK:
        now = time.time()
        cached = _model_cache.get(ck)
        if cached and (now - cached[1]) < _MODEL_CACHE_TTL:
            return cached[0]
        genai.configure(api_key=api_key)
        valid_models = [m.name for m in genai.list_models()
                       if 'generateContent' in m.supported_generation_methods]
        if valid_models:
            # 상한 초과 시 가장 오래된 항목 제거
            if len(_model_cache) >= _MODEL_CACHE_MAX:
                try:
                    oldest = min(_model_cache.items(), key=lambda kv: kv[1][1])[0]
                    _model_cache.pop(oldest, None)
                except Exception:
                    _model_cache.clear()
            _model_cache[ck] = (valid_models, now)
        return valid_models


def _resolve_model_for_choice(api_key, choice, manual_name=""):
    """사용자 선택에 따라 실제 모델명 결정.
    choice: "flash_auto" | "manual"
    Returns: (model_name, display_name) or (None, error_msg)

    v3.0: Pro 경로 완전 제거. 유료 전용 모델은 수동 선택으로도 통과 불가.
    """
    valid = _get_cached_models(api_key)
    if not valid:
        return None, "사용 가능한 Gemini 모델을 찾지 못했습니다."

    if choice == "manual":
        if not manual_name:
            return None, "수동 모드인데 모델명이 비어있습니다."
        # v3.0: 사용자가 어떤 경로로 유료 모델명을 넣더라도 여기서 차단
        if _is_blocked_model(manual_name):
            return None, (
                f"'{manual_name}'은(는) 유료 결제가 필요한 모델이라 사용할 수 없습니다. "
                "이 앱은 무료 티어 Flash 계열만 지원합니다."
            )
        for m in valid:
            if m == manual_name or m.endswith(f"/{manual_name}") or m.split("/")[-1] == manual_name:
                if _is_blocked_model(m):
                    return None, (
                        f"'{manual_name}'은(는) 유료 전용 모델입니다. "
                        "무료 Flash 계열을 선택해주세요."
                    )
                return m, _model_label(m)
        return None, f"입력한 모델 '{manual_name}'을 찾지 못했습니다."

    # flash_auto (기본이자 유일한 자동 모드) — 가장 최신 Flash 자동 선택
    target = _pick_latest_flash(valid)
    if not target:
        return None, (
            "사용 가능한 무료 Flash 모델을 찾지 못했습니다. "
            "API 키가 유효한지 확인해주세요."
        )
    return target, _model_label(target)


def normalize_company(name):
    if not name or not name.strip():
        return name or ""
    cleaned = name.strip()
    lower = cleaned.lower()
    if lower in COMPANY_ALIASES:
        return COMPANY_ALIASES[lower]
    for alias_key, alias_val in COMPANY_ALIASES.items():
        if len(alias_key) >= 3 and alias_key in lower:
            return alias_val
    return cleaned


def _safe_filename(text, max_len=40):
    if not text:
        return "unknown"
    safe = re.sub(r'[\\/:*?"<>|\x00-\x1f\x7f]', '_', str(text))
    safe = re.sub(r'\s+', '_', safe)
    safe = safe.strip('_.')
    WINDOWS_RESERVED = {"CON", "PRN", "AUX", "NUL", "COM1", "COM2", "LPT1", "LPT2"}
    if safe.upper() in WINDOWS_RESERVED:
        safe = f"_{safe}"
    return safe[:max_len] if safe else "unknown"


# ==========================================
# 유틸리티 함수들
# ==========================================
def read_excel_from_bytes(uploaded_file):
    """Excel에서 문서 목록 추출.

    v2.6 Fix C: try/finally로 wb.close() 보장 (반복 업로드 시 메모리 누수 방지).
    v3.3:
      - 하이퍼링크 도메인 검증 (3gpp.org 외 차단)
      - 하이퍼링크 없을 때 RAN1 122차로 고정 조립하던 것 제거.
        (RAN3 문서를 올리면 전부 404가 나던 문제)
        대신 문서번호 접두어로 WG를 추론하고, 추론 불가 시 그 행은 제외.
    Returns: (entries, skipped) — skipped는 [(docid, 사유), ...]
    """
    wb = load_workbook(uploaded_file, read_only=False, data_only=True)
    try:
        ws = wb.active
        entries = []
        skipped = []
        for row in ws.iter_rows(min_row=2):
            cell = row[0]
            comp = row[2] if len(row) > 2 else None
            docid = str(cell.value).strip() if cell.value else ""
            company = normalize_company(str(comp.value).strip()) if comp and comp.value else ""
            if not docid:
                continue

            link = None
            if getattr(cell, "hyperlink", None) and cell.hyperlink.target:
                candidate = str(cell.hyperlink.target).strip()
                if _is_allowed_doc_url(candidate):
                    link = candidate
                else:
                    skipped.append((docid, "허용되지 않은 도메인 링크"))
                    continue

            if not link:
                # v3.3: 하이퍼링크가 없으면 이 앱은 경로를 확정할 수 없음.
                # 예전처럼 RAN1 122차로 찍으면 다른 WG 문서가 전부 실패하므로,
                # 링크 없는 행은 명시적으로 제외하고 사용자에게 알린다.
                skipped.append((docid, "하이퍼링크 없음 (회의 번호 조회 방식 이용 권장)"))
                continue

            entries.append({"doc": docid, "company": company, "link": link})
        return entries, skipped
    finally:
        try:
            wb.close()
        except Exception:
            pass


def _normalize_bis(s):
    if s.startswith("bis"):
        return "bis" + s[3:]
    if s.startswith("b") and (len(s) == 1 or not s[1].isalpha()):
        return "bis" + s[1:]
    return s


WG_FTP_MAP = {
    "RAN1": "tsg_ran/WG1_RL1", "RAN2": "tsg_ran/WG2_RL2",
    "RAN3": "tsg_ran/WG3_Iu", "RAN4": "tsg_ran/WG4_Radio",
    "SA1": "tsg_sa/WG1_Serv", "SA2": "tsg_sa/WG2_Arch",
    "SA3": "tsg_sa/WG3_Security", "SA4": "tsg_sa/WG4_CODEC",
    "SA5": "tsg_sa/WG5_TM", "SA6": "tsg_sa/WG6_MissionCritical",
    "CT1": "tsg_ct/WG1_mm-cc-sm_ex-CN1",
    "CT3": "tsg_ct/WG3_interworking_ex-CN3",
    "CT4": "tsg_ct/WG4_protocollars_ex-CN4",
}

WG_MEETING_PREFIXES = {
    "RAN1": ["TSGR1_"], "RAN2": ["TSGR2_"], "RAN3": ["TSGR3_"], "RAN4": ["TSGR4_"],
    "SA1":  ["TSGS1_"], "SA2":  ["TSGS2_"], "SA3":  ["TSGS3_"], "SA4":  ["TSGS4_"],
    "SA5":  ["TSGS5_"], "SA6":  ["TSGS6_"],
    "CT1":  ["TSGC1_"], "CT3":  ["TSGC3_"], "CT4":  ["CT4_"],
}

WG_TDOC_PREFIX = {
    "RAN1": "TDoc_List_Meeting_RAN1#", "RAN2": "TDoc_List_Meeting_RAN2#",
    "RAN3": "TDoc_List_Meeting_RAN3#", "RAN4": "TDoc_List_Meeting_RAN4#",
    "SA1": "TDoc_List_Meeting_SA1#", "SA2": "TDoc_List_Meeting_SA2#",
    "SA3": "TDoc_List_Meeting_SA3#", "SA4": "TDoc_List_Meeting_SA4#",
    "SA5": "TDoc_List_Meeting_SA5#", "SA6": "TDoc_List_Meeting_SA6#",
    "CT1": "TDoc_List_Meeting_CT1#", "CT3": "TDoc_List_Meeting_CT3#",
    "CT4": "TDoc_List_Meeting_CT4#",
}


def _request_with_retry(url, method="get", max_retries=3, timeout=60, **kwargs):
    # v3.3: verify를 강제하지 않음 → _http_request가 검증 ON으로 시도 후
    #       SSLError일 때만 완화 폴백
    kwargs.setdefault("headers", {"User-Agent": "Mozilla/5.0"})
    kwargs["timeout"] = timeout

    last_error = None
    for attempt in range(max_retries):
        try:
            if method == "head":
                r = _http_head(url, **kwargs)
            else:
                r = _http_get(url, **kwargs)
            if r.status_code == 200:
                return r
            last_error = f"HTTP {r.status_code}"
        except requests.exceptions.Timeout:
            last_error = f"Timeout ({timeout}초)"
            append_log(f"3GPP 서버 타임아웃 (시도 {attempt+1}/{max_retries}): {url[:80]}")
        except requests.exceptions.ConnectionError:
            last_error = "연결 실패"
            append_log(f"3GPP 서버 연결 실패 (시도 {attempt+1}/{max_retries}): {url[:80]}")
        except Exception as e:
            last_error = str(e)

        if attempt < max_retries - 1:
            time.sleep(3 * (attempt + 1))

    append_log(f"3GPP 서버 요청 최종 실패: {last_error}")
    return None


def list_meetings_from_ftp(wg):
    """WG의 3GPP FTP에서 회의 폴더 목록을 조회.

    NOTE (v3.5 점검): 현재 UI 어디에서도 호출되지 않는 미사용 함수.
    회의 선택은 resolve_meeting_folder(회의 번호 직접 입력) 경로를 쓰고 있음.
    기능 제거 원칙에 따라 삭제하지 않고 보존 — 추후 '회의 목록에서 고르기'
    UI를 붙일 때 그대로 재사용 가능.
    """
    ftp_path = WG_FTP_MAP.get(wg)
    if not ftp_path:
        return []
    url = f"https://www.3gpp.org/ftp/{ftp_path}/"
    try:
        r = _http_get(url, timeout=15)
        r.raise_for_status()
        all_links = re.findall(r'href="([^"]*)"', r.text)
        prefixes = WG_MEETING_PREFIXES.get(wg, [])
        meetings = []
        seen = set()
        for link in all_links:
            name = link.rstrip("/").split("/")[-1]
            if not name or name in seen:
                continue
            for pfx in prefixes:
                if name.upper().startswith(pfx.upper()):
                    meetings.append(name)
                    seen.add(name)
                    break
        def sort_key(m):
            nums = re.findall(r'\d+', m)
            return int(nums[0]) if nums else 0
        meetings.sort(key=sort_key, reverse=True)
        return meetings[:30]
    except Exception as e:
        append_log(f"FTP 회의 목록 조회 오류: {e}")
        return []


def resolve_meeting_folder(wg, meeting_num):
    ftp_path = WG_FTP_MAP.get(wg, "")
    prefixes = WG_MEETING_PREFIXES.get(wg, [])
    if not prefixes:
        return None

    base = f"{prefixes[0]}{meeting_num}"
    candidates_to_try = [base]

    num_match = re.match(r'(\d+)', meeting_num)
    if num_match:
        num_part = num_match.group(1)
        suffix = meeting_num[len(num_part):].lower().lstrip("-_")

        if suffix:
            all_suffixes = set()
            all_suffixes.add(suffix)
            all_suffixes.add(f"-{suffix}")
            all_suffixes.add(f"_{suffix}")
            if suffix == "bis":
                all_suffixes.add("b")
                all_suffixes.add("-b")
            elif suffix == "b":
                all_suffixes.add("bis")
                all_suffixes.add("-bis")

            for sfx in all_suffixes:
                candidate = f"{prefixes[0]}{num_part}{sfx}"
                if candidate not in candidates_to_try:
                    candidates_to_try.append(candidate)

    def _try_head(candidate):
        test_url = f"https://www.3gpp.org/ftp/{ftp_path}/{candidate}/Docs/"
        r = _request_with_retry(test_url, method="head", max_retries=2, timeout=15)
        return candidate, (r and r.status_code == 200)

    if len(candidates_to_try) == 1:
        cand, ok = _try_head(candidates_to_try[0])
        if ok:
            return cand
    else:
        results = {}
        with ThreadPoolExecutor(max_workers=min(len(candidates_to_try), 5)) as ex:
            futures = {ex.submit(_try_head, c): c for c in candidates_to_try}
            for fut in as_completed(futures):
                try:
                    cand, ok = fut.result()
                    results[cand] = ok
                except Exception:
                    results[futures[fut]] = False
        for c in candidates_to_try:
            if results.get(c):
                return c

    dir_url = f"https://www.3gpp.org/ftp/{ftp_path}/"
    r = _request_with_retry(dir_url, max_retries=3, timeout=30)
    if not r:
        return None
    try:
        all_links = re.findall(r'href="([^"?]+)"', r.text)
        all_links += re.findall(r'>([A-Z][^<]{3,80})<', r.text)

        num_match = re.match(r'(\d+)', meeting_num)
        if not num_match:
            return None
        num_part = num_match.group(1)
        suffix_part = meeting_num[len(num_part):].lower()

        found = []
        for link in all_links:
            name = link.rstrip("/").split("/")[-1].strip()
            if not name:
                continue
            name_upper = name.upper()
            prefix_upper = prefixes[0].upper()
            if not name_upper.startswith(prefix_upper):
                continue
            after_prefix = name[len(prefixes[0]):]
            folder_num_match = re.match(r'(\d+)', after_prefix)
            if not folder_num_match:
                continue
            folder_num = folder_num_match.group(1)
            if folder_num != num_part:
                continue
            folder_rest = after_prefix[len(folder_num):].lower().lstrip("-_")
            if suffix_part:
                normalized_suffix = suffix_part.replace("-", "").replace("_", "")
                normalized_folder = folder_rest.replace("-", "").replace("_", "")
                if _normalize_bis(normalized_folder).startswith(_normalize_bis(normalized_suffix)):
                    found.append(name)
            else:
                if (folder_rest == "" or
                    folder_rest.startswith("_") or
                    folder_rest.startswith("/") or
                    re.match(r'^(e|ah|ahe?)[\W_]', folder_rest, re.I) or
                    re.match(r'^(e|ah|ahe?)$', folder_rest, re.I)):
                    found.append(name)

        if found:
            found.sort(key=len)
            append_log(f"폴더 후보: {found}")
            return found[0]

    except Exception as e:
        append_log(f"폴더 검색 오류: {e}")

    return None


def fetch_tdoc_list_xlsx(wg, meeting_folder):
    import urllib.parse

    ftp_path = WG_FTP_MAP.get(wg, "")
    tdoc_prefix = WG_TDOC_PREFIX.get(wg, "TDoc_List_Meeting_")

    meeting_num = meeting_folder
    for pfx in WG_MEETING_PREFIXES.get(wg, []):
        if meeting_folder.upper().startswith(pfx.upper()):
            meeting_num = meeting_folder[len(pfx):]
            break

    match = re.match(r'^(\d+(?:-?bis|-?e|-?b)?)', meeting_num, re.I)
    if match:
        meeting_num = match.group(1)

    docs_url = f"https://www.3gpp.org/ftp/{ftp_path}/{meeting_folder}/Docs/"

    xlsx_candidates = [f"{tdoc_prefix}{meeting_num}.xlsx"]

    mn_match = re.match(r'^(\d+)[-_]?(.*)', meeting_num, re.I)
    if mn_match:
        mn_num = mn_match.group(1)
        mn_suffix = mn_match.group(2).lower()

        suffix_variants = set()
        if mn_suffix:
            suffix_variants.add(mn_suffix)
            suffix_variants.add(f"-{mn_suffix}")
            if mn_suffix == "b":
                suffix_variants.add("bis")
                suffix_variants.add("-bis")
            elif mn_suffix == "bis":
                suffix_variants.add("b")
                suffix_variants.add("-b")

            for sfx in suffix_variants:
                candidate = f"{tdoc_prefix}{mn_num}{sfx}.xlsx"
                if candidate not in xlsx_candidates:
                    xlsx_candidates.append(candidate)

    r = None
    for xlsx_filename in xlsx_candidates:
        xlsx_url_encoded = f"{docs_url}{urllib.parse.quote(xlsx_filename)}"
        resp = _request_with_retry(xlsx_url_encoded, max_retries=3, timeout=60)
        if resp and resp.status_code == 200:
            r = resp
            append_log(f"TDoc xlsx 발견: {xlsx_filename}")
            break

    if r is None:
        dir_resp = _request_with_retry(docs_url, max_retries=2, timeout=30)
        if dir_resp:
            try:
                xlsx_links = re.findall(r'href="([^"]*TDoc_List[^"]*\.xlsx)"', dir_resp.text, re.I)
                if xlsx_links:
                    actual_filename = xlsx_links[0].split("/")[-1]
                    actual_url = f"{docs_url}{urllib.parse.quote(actual_filename)}"
                    r = _request_with_retry(actual_url, max_retries=3, timeout=60)
            except Exception as e:
                append_log(f"TDoc 리스트 디렉토리 검색 실패: {e}")

    if r is None:
        append_log(f"TDoc 리스트 다운로드 최종 실패")
        return {}, []

    wb = load_workbook(io.BytesIO(r.content), read_only=True, data_only=True)
    ws = wb.active

    header_row = None
    col_map = {}
    for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=10), start=1):
        for col_idx, cell in enumerate(row):
            val = str(cell.value or "").strip().lower()
            if any(kw in val for kw in ["tdoc", "td#", "td number"]) and "tdoc" not in col_map:
                col_map["tdoc"] = col_idx
            if any(kw in val for kw in ["source", "company", "submitting"]) and "company" not in col_map:
                col_map["company"] = col_idx
            if "agenda" in val:
                if "description" in val:
                    col_map["agenda"] = col_idx
                elif "agenda" not in col_map:
                    col_map["agenda_num"] = col_idx
        if "tdoc" in col_map and ("agenda" in col_map or "agenda_num" in col_map):
            header_row = row_idx
            break

    if "agenda" not in col_map and "agenda_num" in col_map:
        col_map["agenda"] = col_map["agenda_num"]

    has_separate_num = "agenda_num" in col_map and "agenda" in col_map and col_map.get("agenda_num") != col_map.get("agenda")

    if not header_row:
        header_row = 1
        col_map = {"tdoc": 0, "company": 2, "agenda_num": 10, "agenda": 11}
        has_separate_num = True

    entries = []
    agenda_dict = {}

    for row in ws.iter_rows(min_row=header_row + 1):
        tdoc_idx = col_map.get("tdoc", 0)
        company_idx = col_map.get("company", 2)
        agenda_idx = col_map.get("agenda", 11)
        agenda_num_idx = col_map.get("agenda_num", 10)

        if len(row) <= tdoc_idx:
            continue

        tdoc_cell = row[tdoc_idx]
        company_cell = row[company_idx] if len(row) > company_idx else None
        agenda_cell = row[agenda_idx] if len(row) > agenda_idx else None
        agenda_num_cell = row[agenda_num_idx] if has_separate_num and len(row) > agenda_num_idx else None

        tdoc_id = str(tdoc_cell.value or "").strip()
        if not tdoc_id:
            continue

        company = normalize_company(str(company_cell.value or "").strip()) if company_cell else ""
        agenda_desc = str(agenda_cell.value or "").strip() if agenda_cell else ""
        agenda_num = str(agenda_num_cell.value or "").strip() if agenda_num_cell else ""

        if agenda_num and agenda_desc and agenda_num != agenda_desc:
            agenda = f"{agenda_num} - {agenda_desc}"
        elif agenda_desc:
            agenda = agenda_desc
        elif agenda_num:
            agenda = agenda_num
        else:
            agenda = ""

        if getattr(tdoc_cell, "hyperlink", None) and tdoc_cell.hyperlink.target:
            link = tdoc_cell.hyperlink.target
        else:
            link = f"{docs_url}{tdoc_id}.zip"

        entry = {"doc": tdoc_id, "company": company, "link": link, "agenda": agenda}
        entries.append(entry)

        if agenda:
            agenda_dict.setdefault(agenda, [])
            agenda_dict[agenda].append(entry)

    wb.close()
    return agenda_dict, entries


def clone_paragraph(src, dest):
    np_para = dest.add_paragraph("", style=src.style)
    for r in src.runs:
        try:
            nr = np_para.add_run(_xml_safe(r.text))
            nr.bold = r.bold
            nr.italic = r.italic
            nr.underline = r.underline
            if hasattr(r.font, "name") and r.font.name:
                nr.font.name = r.font.name
            if hasattr(r.font, "size") and r.font.size:
                nr.font.size = r.font.size
            if hasattr(r.font, "color") and getattr(r.font.color, "rgb", None):
                nr.font.color.rgb = r.font.color.rgb
        except Exception:
            try:
                np_para.add_run(_xml_safe(r.text))
            except Exception:
                pass
    return np_para


# ==========================================
# ★ v2.6 Fix F: ZipSlip 방어 (os.path.commonpath 기반) ★
# ==========================================
def _safe_extractall(zf, target_dir):
    """zip 멤버가 target_dir 밖으로 escape하지 않을 때만 extractall.
    악성 멤버 발견 시 즉시 ValueError 발생, 부분 추출 안 함.

    naive startswith 비교 대신 os.path.commonpath 사용:
    - symlink resolution 정확
    - Unicode normalization 정확
    - Windows 드라이브 차이 정확

    v3.3 추가: 압축 폭탄(zip bomb) 방어.
    Cloud Run에서 /tmp가 RAM이므로, 작은 zip이 거대하게 풀리면
    인스턴스가 OOM으로 죽음. 해제 후 총 크기와 압축률을 사전 검사.
    """
    target_real = os.path.realpath(target_dir)

    # v3.3: 압축 해제 총 크기 / 압축률 사전 검사 (헤더만 읽으므로 비용 거의 없음)
    try:
        infos = zf.infolist()
        total_uncompressed = sum(getattr(i, "file_size", 0) or 0 for i in infos)
        total_compressed = sum(getattr(i, "compress_size", 0) or 0 for i in infos)
        if total_uncompressed > MAX_EXTRACT_BYTES:
            raise ValueError(
                f"압축 해제 크기 초과 ({total_uncompressed / 1048576:.0f}MB > "
                f"{MAX_EXTRACT_BYTES // 1048576}MB)"
            )
        if total_compressed > 0:
            ratio = total_uncompressed / total_compressed
            if ratio > MAX_COMPRESSION_RATIO:
                raise ValueError(
                    f"비정상 압축률 감지 ({ratio:.0f}:1) — 압축 폭탄 가능성"
                )
    except ValueError:
        raise
    except Exception:
        # infolist 읽기 실패는 아래 경로 검사에서 다시 걸림
        pass

    for member in zf.namelist():
        if not member or os.path.isabs(member):
            raise ValueError(f"ZipSlip vulnerability detected (absolute path): {member}")
        member_path = os.path.realpath(os.path.join(target_real, member))
        try:
            common = os.path.commonpath([target_real, member_path])
        except ValueError:
            # commonpath는 빈 입력 또는 mixed-drive (Windows)에서 ValueError 발생
            # 둘 다 boundary 위반으로 간주
            raise ValueError(f"ZipSlip vulnerability detected (path mismatch): {member}")
        if common != target_real:
            raise ValueError(f"ZipSlip vulnerability detected: {member}")
    zf.extractall(target_dir)


def repackage_docm_to_docx(path, td):
    """v2.6 Fix D: 작업 완료 후 docm_unzip 폴더 즉시 정리 (분석 중 디스크 누적 방지).
    v2.6 Fix F: _safe_extractall로 ZipSlip 방어."""
    ud = os.path.join(td, "docm_unzip")
    os.makedirs(ud, exist_ok=True)
    with zipfile.ZipFile(path, 'r') as zf:
        _safe_extractall(zf, ud)
    tf = os.path.join(ud, "[Content_Types].xml")
    if not os.path.exists(tf):
        return path
    with open(tf, 'r', encoding='utf-8') as f:
        t = f.read()
    t = t.replace(
        'application/vnd.ms-word.document.macroEnabled.main+xml',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml'
    )
    with open(tf, 'w', encoding='utf-8') as f:
        f.write(t)
    rp = os.path.join(td, "repack.zip")
    with zipfile.ZipFile(rp, 'w', zipfile.ZIP_DEFLATED) as zf:
        for r, _, fs in os.walk(ud):
            for f in fs:
                full = os.path.join(r, f)
                arc = os.path.relpath(full, ud)
                zf.write(full, arc)
    out = os.path.join(td, "repack.docx")
    os.rename(rp, out)
    # v2.6 Fix D: 중간 작업물 즉시 정리
    try:
        shutil.rmtree(ud, ignore_errors=True)
    except Exception:
        pass
    return out


def _download_doc(entry, td_name, headers, max_retries=3):
    """문서 zip 다운로드.

    v3.3 추가 방어:
    1) 도메인 화이트리스트 — 3gpp.org 외에는 요청 자체를 안 함 (최종 방어선)
    2) 외부 리다이렉트 차단 — 3gpp.org 밖으로 나가는 리다이렉트는 거부
    3) 파일당 크기 상한 (MAX_DOWNLOAD_BYTES) — Content-Length 선검사 +
       스트리밍 중 실시간 검사 (Content-Length 없는 서버 대비)
    4) 실행당 누적 상한 (MAX_TOTAL_DOWNLOAD_BYTES)
    5) stream=True — r.content로 전체를 메모리에 올리지 않음.
       Cloud Run에서 /tmp가 RAM이므로 메모리 2배 사용을 피하는 효과.
    """
    url = entry.get("link", "")

    # (1) 최종 방어선: 허용 도메인이 아니면 요청조차 하지 않음
    if not _is_allowed_doc_url(url):
        return entry, None, "허용되지 않은 도메인 (3gpp.org만 허용)"

    last_error = None
    for attempt in range(max_retries):
        fp = None
        try:
            cur_url = url
            r = None
            # (2) 리다이렉트를 수동 추적하며 도메인 이탈 차단 (최대 5회)
            for _ in range(5):
                r = _http_get(cur_url, headers=headers, timeout=60,
                              stream=True, allow_redirects=False)
                if r.status_code in (301, 302, 303, 307, 308):
                    loc = r.headers.get("Location", "")
                    nxt = urllib.parse.urljoin(cur_url, loc)
                    try:
                        r.close()
                    except Exception:
                        pass
                    if not _is_allowed_doc_url(nxt):
                        return entry, None, "허용되지 않은 도메인으로 리다이렉트"
                    cur_url = nxt
                    continue
                break
            if r is None:
                raise requests.exceptions.ConnectionError("no response")
            r.raise_for_status()

            # (3) Content-Length 선검사
            try:
                clen = int(r.headers.get("Content-Length") or 0)
            except Exception:
                clen = 0
            if clen and clen > MAX_DOWNLOAD_BYTES:
                try:
                    r.close()
                except Exception:
                    pass
                return entry, None, (
                    f"파일이 너무 큼 ({clen / 1048576:.1f}MB > "
                    f"{MAX_DOWNLOAD_BYTES // 1048576}MB 제한)"
                )

            # (5) 스트리밍 저장 + 실시간 크기 검사
            fp = os.path.join(td_name, f"{entry['doc']}.zip")
            written = 0
            with open(fp, "wb") as f:
                for chunk in r.iter_content(chunk_size=1 << 16):
                    if not chunk:
                        continue
                    written += len(chunk)
                    if written > MAX_DOWNLOAD_BYTES:
                        raise ValueError(
                            f"다운로드 중 크기 초과 "
                            f"({MAX_DOWNLOAD_BYTES // 1048576}MB 제한)"
                        )
                    f.write(chunk)
            try:
                r.close()
            except Exception:
                pass

            # (4) 실행당 누적 상한
            if not _add_download_bytes(written):
                _safe_unlink(fp)
                return entry, None, (
                    f"실행당 누적 다운로드 상한 초과 "
                    f"({MAX_TOTAL_DOWNLOAD_BYTES // 1048576}MB)"
                )

            return entry, fp, None

        except (requests.exceptions.Timeout, requests.exceptions.ConnectionError) as ex:
            _safe_unlink(fp)
            last_error = str(ex)
            if attempt < max_retries - 1:
                time.sleep(2 * (attempt + 1))
        except ValueError as ex:
            # 크기 초과 등 — 재시도해도 같으므로 즉시 종료
            _safe_unlink(fp)
            return entry, None, str(ex)
        except Exception as ex:
            _safe_unlink(fp)
            return entry, None, str(ex)
    return entry, None, last_error or "Download failed after retries"


def _safe_unlink(path):
    """부분 다운로드 파일 등을 조용히 삭제."""
    try:
        if path and os.path.isfile(path):
            os.remove(path)
    except Exception:
        pass


def _release_doc_workspace(zip_path, extract_dir):
    """v2.9: 한 문서 처리가 끝나면 그 문서의 zip과 압축 해제 폴더를 즉시 해제.

    Cloud Run에서 /tmp는 메모리(tmpfs)이므로, 이걸 안 하면 200개 문서의
    zip + 압축 결과가 전부 RAM에 누적되어 OOM 위험. 문서 단위로 즉시
    비우면 피크 메모리가 '전체 합계'에서 '동시 처리분'으로 내려감.
    이미 텍스트는 extracted_list에 추출된 뒤이므로 기능 영향 없음.
    """
    try:
        if extract_dir and os.path.isdir(extract_dir):
            shutil.rmtree(extract_dir, ignore_errors=True)
    except Exception:
        pass
    try:
        if zip_path and os.path.isfile(zip_path):
            os.remove(zip_path)
    except Exception:
        pass


# ==========================================
# v2.8: 디스크 청소 — 네임스페이스 안전 버전
# 기존(~v2.6)은 ^tmp[a-zA-Z0-9_]{6,}$ 패턴으로 /tmp 전역을 훑어
# 다른 프로세스의 임시 디렉토리까지 삭제 대상으로 잡는 위험이 있었음.
# v2.8은 우리 네임스페이스(/tmp/3gpp_analyzer/run_*)만 PID/heartbeat로
# 판단해 회수. 남의 파일은 어떤 경우에도 안 건드림.
# ==========================================
def _cleanup_tmp_if_low_disk(force=False):
    """네임스페이스 내 stale run 디렉토리만 안전하게 회수.
    (force 인자는 하위 호환을 위해 유지하나, 동작은 항상 안전한 reaper)."""
    try:
        reaped = reap_stale_run_dirs()
        if reaped:
            append_log(f"🧹 정리: 중단된 작업 잔해 {reaped}개 회수")
    except Exception:
        pass


def extract_all_conclusions(entries, status_elem, progress_elem, log_func):
    _cleanup_tmp_if_low_disk()
    if CLOUD_FUNCTION_URL:
        return _extract_via_cloud(entries, status_elem, progress_elem, log_func)
    return _extract_local(entries, status_elem, progress_elem, log_func)


def _extract_via_cloud(entries, status_elem, progress_elem, log_func):
    od = Document()
    od.add_heading("3GPP Conclusions", level=0)
    extracted_list = []
    total = len(entries)
    if total == 0:
        log_func("입력 문서 없음")
        bio = io.BytesIO()
        od.save(bio)
        bio.seek(0)
        return bio
    batch_size = 20
    all_results = []

    for i in range(0, total, batch_size):
        batch = entries[i:i + batch_size]
        status_elem.text(f"☁️ 클라우드 처리 [{min(i+batch_size, total)}/{total}]")
        progress_elem.progress(min(i+batch_size, total) / max(total, 1))
        try:
            resp = requests.post(CLOUD_FUNCTION_URL, json={"entries": batch}, timeout=300)
            resp.raise_for_status()
            all_results.extend(resp.json().get("results", []))
        except Exception as e:
            log_func(f"Cloud Function 오류, 로컬 전환: {e}")
            return _extract_local(entries, status_elem, progress_elem, log_func)

    for idx, item in enumerate(all_results, 1):
        try:
            tbl = od.add_table(rows=4, cols=2, style="Table Grid")
            _safe_set_cell_text(tbl.cell(0, 0), "Document")
            _safe_set_cell_text(tbl.cell(0, 1), item.get("doc", ""))
            _safe_set_cell_text(tbl.cell(1, 0), "Link")
            _safe_set_cell_text(tbl.cell(1, 1), item.get("link", ""))
            _safe_set_cell_text(tbl.cell(2, 0), "Company")
            _safe_set_cell_text(tbl.cell(2, 1), item.get("company", ""))
            _safe_set_cell_text(tbl.cell(3, 0), "Title")
            _safe_set_cell_text(tbl.cell(3, 1), item.get("title", ""))

            content = item.get("content", "") or ""
            is_cr_cloud = bool(content) and (
                content.lstrip().startswith("[CR — Change Request 문서]") or
                content.lstrip().startswith("[CR \u2014 Change Request")
            )

            if content and content not in ("결론 섹션 없음", "DOC 파일 없음"):
                for line in content.split("\n"):
                    if line.strip():
                        _safe_add_paragraph(od, line)
            else:
                _safe_add_paragraph(od, content or "결론 섹션 없음")

            extracted_list.append({
                "doc": item.get("doc", ""), "company": item.get("company", ""),
                "link": item.get("link", ""), "title": item.get("title", ""),
                "is_cr": is_cr_cloud,
                "content": content,
                "full_content": content,
            })
            log_func(f"{item.get('doc','')} 추출 완료")
        except Exception as ex:
            log_func(f"Cloud 아이템 처리 오류 ({item.get('doc','?')}): {ex}")
            try:
                _safe_add_paragraph(od, f"[처리 오류] {item.get('doc','?')}: {str(ex)[:200]}")
            except Exception:
                pass
            extracted_list.append({
                "doc": item.get("doc", ""), "company": item.get("company", ""),
                "link": item.get("link", ""), "title": "(처리 오류)",
                "is_cr": False,
                "content": f"처리 오류: {str(ex)[:200]}",
                "full_content": "",
            })

        if idx < len(all_results):
            try:
                od.add_page_break()
            except Exception:
                pass

    st.session_state.extracted_data = extracted_list
    _build_notebooklm_txt(extracted_list)
    bio = io.BytesIO()
    od.save(bio)
    bio.seek(0)
    return bio


def _extract_local(entries, status_elem, progress_elem, log_func):
    with _managed_run_dir() as temp_dir:
        log_func("임시 작업 디렉터리 생성 (자동 정리)")
        _reset_download_budget()   # v3.3: 실행마다 누적 다운로드 예산 초기화
        od = Document()
        od.add_heading("3GPP Conclusions", level=0)

        # v2.6 Fix J: CONCLUSION_PATTERNS/END_PATTERNS 변수 제거.
        # 헤더 판정은 is_conclusion_header() / is_end_header() 함수 직접 호출.
        headers = {"User-Agent": "Mozilla/5.0"}
        download_results = []
        extracted_list = []
        total = len(entries)

        if total == 0:
            log_func("입력 문서 없음")
            bio = io.BytesIO()
            od.save(bio)
            bio.seek(0)
            return bio

        with ThreadPoolExecutor(max_workers=10) as executor:
            futures = {executor.submit(_download_doc, e, temp_dir, headers): e for e in entries}
            for i, fut in enumerate(as_completed(futures), start=1):
                e, fp, err = fut.result()
                download_results.append((e, fp, err))
                progress_elem.progress(i / max(total, 1))
                status_elem.text(f"Downloaded [{i}/{total}]: {e['doc']}")
                log_func(f"[{i}/{total}] Downloaded: {e['doc']}")
                if i % 10 == 0:
                    _touch_heartbeat(temp_dir)  # v2.8: 장시간 작업 중 heartbeat 갱신

        for idx, (e, fp, err) in enumerate(download_results, start=1):
            status_elem.text(f"Extracting [{idx}/{total}]: {e['doc']}")
            if idx % 10 == 0:
                _touch_heartbeat(temp_dir)  # v2.8: 추출 중 heartbeat 갱신
            doc_text_buffer = []
            full_text_buffer = []

            tbl = od.add_table(rows=4, cols=2, style="Table Grid")
            _safe_set_cell_text(tbl.cell(0, 0), "Document")
            _safe_set_cell_text(tbl.cell(0, 1), e["doc"])
            _safe_set_cell_text(tbl.cell(1, 0), "Link")
            _safe_set_cell_text(tbl.cell(1, 1), e["link"])
            _safe_set_cell_text(tbl.cell(2, 0), "Company")
            _safe_set_cell_text(tbl.cell(2, 1), e["company"])
            _safe_set_cell_text(tbl.cell(3, 0), "Title")

            # v3.3: 추출 경로를 한 곳에서만 계산해 finally와 공유.
            # (v3.2까지는 try 안과 finally에서 각각 조립해, 한쪽만 바꾸면
            #  정리가 조용히 실패하는 구조였음)
            # v2.9: idx를 붙여 동일 문서번호 중복 시 서로 덮어쓰지 않도록
            ed = os.path.join(temp_dir, f"{e['doc']}__{idx}")

            try:
                if err or not fp:
                    raise Exception(err or "Download failed")
                os.makedirs(ed, exist_ok=True)
                with zipfile.ZipFile(fp) as zf:
                    # v2.6 Fix F: ZipSlip 방어
                    _safe_extractall(zf, ed)

                # v2.6 Fix G: macOS AppleDouble 아티팩트(._filename.docx 등) 제외.
                # 이런 파일을 python-docx가 열려고 하면 BadZipFile로 실패함.
                src_path = None
                for ext in ("*.docx", "*.docm", "*.doc", "*.pptx", "*.ppt", "*.pdf"):
                    for candidate in Path(ed).rglob(ext):
                        if candidate.name.startswith("._"):
                            continue
                        src_path = candidate
                        break
                    if src_path: break

                if not src_path:
                    _safe_add_paragraph(od, "문서 파일을 찾을 수 없습니다 (docx/doc/ppt/pdf).")
                    log_func(f"{e['doc']} 파일 없음")
                    extracted_list.append({
                        "doc": e["doc"], "company": e["company"], "link": e["link"],
                        "title": "(파일 없음)",
                        "content": "zip 내부에 문서 파일이 없습니다 (docx/doc/ppt/pdf)",
                        "full_content": ""
                    })
                    if idx < len(download_results):
                        od.add_page_break()
                    continue

                file_path_str = str(src_path)

                if src_path.suffix.lower() == ".doc":
                    doc_appended = False
                    try:
                        with open(file_path_str, "rb") as bf:
                            raw = bf.read()
                        text_chunks = re.findall(rb'[\x20-\x7E]{20,}', raw)
                        raw_text = "\n".join(chunk.decode('ascii', errors='ignore') for chunk in text_chunks)
                        if raw_text:
                            _safe_add_paragraph(od, "[구형 .doc — 텍스트 추출 (서식 없음)]")
                            lines = raw_text.split('\n')
                            found_conclusion = False
                            for li, line in enumerate(lines):
                                if re.search(r'(?:conclusion|summary)', line, re.I):
                                    found_conclusion = True
                                    for cl in lines[li:li+30]:
                                        _safe_add_paragraph(od, cl)
                                        doc_text_buffer.append(cl)
                                    break
                            if not found_conclusion:
                                _safe_add_paragraph(od, "결론 섹션 없음 (구형 .doc)")
                                for cl in lines[-20:]:
                                    doc_text_buffer.append(cl)
                            extracted_list.append({
                                "doc": e["doc"], "company": e["company"], "link": e["link"],
                                "title": "(구형 .doc)",
                                "content": "\n".join(doc_text_buffer) if doc_text_buffer else "텍스트 추출 실패",
                                "full_content": raw_text[:5000]
                            })
                            doc_appended = True
                            log_func(f"{e['doc']} .doc 텍스트 추출")
                        else:
                            _safe_add_paragraph(od, "구형 .doc 파일에서 텍스트를 추출할 수 없습니다.")
                            log_func(f"{e['doc']} .doc 텍스트 추출 실패")
                    except Exception as ex:
                        _safe_add_paragraph(od, f"구형 .doc 파일 처리 오류: {ex}")
                        log_func(f"{e['doc']} .doc 오류: {ex}")
                    if not doc_appended:
                        extracted_list.append({
                            "doc": e["doc"], "company": e["company"], "link": e["link"],
                            "title": "(구형 .doc — 추출 실패)",
                            "content": "구형 .doc 텍스트 추출 실패",
                            "full_content": ""
                        })
                    if idx < len(download_results):
                        od.add_page_break()
                    continue

                if src_path.suffix.lower() in (".pptx", ".ppt"):
                    ppt_appended = False
                    try:
                        from pptx import Presentation
                        prs = Presentation(file_path_str)
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    texts.append(shape.text.strip())
                        ppt_text = "\n".join(texts)
                        _safe_add_paragraph(od, "[PPT 문서 — 슬라이드 텍스트 추출]")
                        for line in ppt_text.split('\n')[:50]:
                            if line.strip():
                                _safe_add_paragraph(od, line.strip())
                                doc_text_buffer.append(line.strip())
                        extracted_list.append({
                            "doc": e["doc"], "company": e["company"], "link": e["link"],
                            "title": "(PPT)", "content": ppt_text[:5000],
                            "full_content": ppt_text
                        })
                        ppt_appended = True
                        log_func(f"{e['doc']} PPT 추출 완료")
                    except ImportError:
                        _safe_add_paragraph(od, "PPT 파싱 라이브러리 없음 (python-pptx)")
                    except Exception as ex:
                        _safe_add_paragraph(od, f"PPT 처리 오류: {ex}")
                        log_func(f"{e['doc']} PPT 오류: {ex}")
                    if not ppt_appended:
                        extracted_list.append({
                            "doc": e["doc"], "company": e["company"], "link": e["link"],
                            "title": "(PPT — 추출 실패)",
                            "content": "PPT 텍스트 추출 실패",
                            "full_content": ""
                        })
                    if idx < len(download_results):
                        od.add_page_break()
                    continue

                if src_path.suffix.lower() == ".pdf":
                    pdf_appended = False
                    try:
                        import fitz
                        pdf_doc = fitz.open(file_path_str)
                        try:
                            pdf_texts = [page.get_text() for page in pdf_doc]
                        finally:
                            pdf_doc.close()
                        pdf_text = "\n".join(pdf_texts)
                        _safe_add_paragraph(od, "[PDF 문서 — 텍스트 추출]")
                        for line in pdf_text.split('\n')[:50]:
                            if line.strip():
                                _safe_add_paragraph(od, line.strip())
                                doc_text_buffer.append(line.strip())
                        extracted_list.append({
                            "doc": e["doc"], "company": e["company"], "link": e["link"],
                            "title": "(PDF)", "content": pdf_text[:5000],
                            "full_content": pdf_text
                        })
                        pdf_appended = True
                        log_func(f"{e['doc']} PDF 추출 완료")
                    except ImportError:
                        # ★ v2.6 Fix E: 바이너리 폴백 제거 ★
                        # 기존 코드는 PDF 바이너리를 ASCII로 강제 디코드해서 그 결과를
                        # extracted_list에 넣었음. PDF는 FlateDecode 압축이라
                        # 의미 없는 PostScript 조각이 LLM 입력으로 들어가
                        # 할루시네이션을 유발. 이제는 명시적 placeholder만 남기고
                        # full_content=""로 두어 Gemini 입력에서 자동 제외.
                        _safe_add_paragraph(od, "[PDF — PyMuPDF 미설치, 텍스트 추출 불가]")
                        extracted_list.append({
                            "doc": e["doc"], "company": e["company"], "link": e["link"],
                            "title": "(PDF — 라이브러리 없음)",
                            "content": "PyMuPDF 라이브러리가 설치되지 않아 PDF 텍스트를 추출할 수 없습니다.",
                            "full_content": ""
                        })
                        pdf_appended = True
                        # 스킵 카운트 증가 → 추출 후 배너 표시용
                        try:
                            st.session_state.pdf_skip_count = st.session_state.get("pdf_skip_count", 0) + 1
                        except Exception:
                            pass
                        log_func(f"{e['doc']} PDF skipped (no PyMuPDF)")
                    except Exception as ex:
                        _safe_add_paragraph(od, f"PDF 처리 오류: {ex}")
                        log_func(f"{e['doc']} PDF 오류: {ex}")
                    if not pdf_appended:
                        extracted_list.append({
                            "doc": e["doc"], "company": e["company"], "link": e["link"],
                            "title": "(PDF — 추출 실패)",
                            "content": "PDF 텍스트 추출 실패",
                            "full_content": ""
                        })
                    if idx < len(download_results):
                        od.add_page_break()
                    continue

                if src_path.suffix.lower() == ".docm":
                    try:
                        file_path_str = repackage_docm_to_docx(file_path_str, temp_dir)
                    except Exception as ex:
                        log_func(f"{e['doc']} docm 변환 오류: {ex}")

                try:
                    sd = Document(file_path_str)
                except Exception as ex:
                    _safe_add_paragraph(od, f"문서를 열 수 없습니다 (구형 .doc 파일이거나 손상됨): {ex}")
                    log_func(f"{e['doc']} 문서 파싱 에러: {ex}")
                    extracted_list.append({
                        "doc": e["doc"], "company": e["company"], "link": e["link"],
                        "title": "(문서 열기 실패)",
                        "content": f"문서를 열 수 없습니다: {str(ex)[:200]}",
                        "full_content": ""
                    })
                    if idx < len(download_results):
                        od.add_page_break()
                    continue

                title = ""
                paras = sd.paragraphs
                for p in paras:
                    t = p.text.strip()
                    if t:
                        full_text_buffer.append(t)
                    if not title and t.lower().startswith("title:"):
                        title = t.split(":", 1)[1].strip()
                if not title:
                    title = sd.core_properties.title or ""

                is_cr = False
                cr_reason = ""
                cr_summary = ""
                cr_title = ""
                try:
                    for tbl_idx, doc_tbl in enumerate(sd.tables[:3]):
                        for row in doc_tbl.rows:
                            row_text = " ".join(cell.text.strip() for cell in row.cells).lower()
                            if "change request" in row_text or "cr-form" in row_text:
                                is_cr = True
                                break
                        if is_cr:
                            break

                    if is_cr:
                        for doc_tbl in sd.tables[:3]:
                            for row in doc_tbl.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                cells_lower = [c.lower() for c in cells]
                                row_joined = " ".join(cells_lower)

                                if "title:" in cells_lower[0] and not cr_title:
                                    for c in cells[1:]:
                                        if c and c != cells[0]:
                                            cr_title = c
                                            break

                                if "reason for change" in row_joined:
                                    for c in cells:
                                        if c.lower() not in ("", "reason for change:", "reason for change"):
                                            cr_reason = c
                                            break

                                if "summary of change" in row_joined:
                                    for c in cells:
                                        if c.lower() not in ("", "summary of change:", "summary of change"):
                                            cr_summary = c
                                            break

                        if cr_title and not title:
                            title = cr_title
                except Exception:
                    pass

                _safe_set_cell_text(tbl.cell(3, 1), title)

                if is_cr:
                    try:
                        hdr_para = _safe_add_paragraph(od, "📋 [CR — Change Request 문서]")
                        if hdr_para.runs:
                            hdr_para.runs[0].bold = True
                    except Exception:
                        pass
                    if cr_reason:
                        try:
                            p_label = od.add_paragraph("")
                            p_label.add_run(_xml_safe("Reason for change: ")).bold = True
                        except Exception:
                            pass
                        _safe_add_paragraph(od, cr_reason)
                        doc_text_buffer.append(f"Reason for change: {cr_reason}")
                    if cr_summary:
                        try:
                            p_label = od.add_paragraph("")
                            p_label.add_run(_xml_safe("Summary of change: ")).bold = True
                        except Exception:
                            pass
                        _safe_add_paragraph(od, cr_summary)
                        doc_text_buffer.append(f"Summary of change: {cr_summary}")
                    if not cr_reason and not cr_summary:
                        _safe_add_paragraph(od, "CR 테이블에서 Reason/Summary를 추출하지 못했습니다.")
                    log_func(f"{e['doc']} CR 문서 추출 완료")

                    extracted_list.append({
                        "doc": e["doc"], "company": e["company"], "link": e["link"],
                        "title": title, "is_cr": True,
                        "content": "\n".join(doc_text_buffer) if doc_text_buffer else "CR 내용 추출 실패",
                        "full_content": ("\n".join(full_text_buffer))[:30000] if full_text_buffer else ""
                    })
                    if idx < len(download_results):
                        od.add_page_break()
                    continue

                # v2.6 Fix J: 어근 기반 헤더 매칭으로 conclusion 섹션 시작 찾기
                start = None
                for j, p in enumerate(paras):
                    if is_conclusion_header(p.text):
                        start = j
                        break

                if start is None:
                    _safe_add_paragraph(od, "결론 섹션 없음")
                    log_func(f"{e['doc']} 결론없음")
                else:
                    # v2.6 Fix J: 어근 기반 헤더 매칭으로 end 섹션 찾기
                    end = len(paras)
                    for j, p in enumerate(paras[start+1:], start+1):
                        if is_end_header(p.text):
                            end = j
                            break
                    for j in range(start+1, end):
                        try:
                            clone_paragraph(paras[j], od)
                        except Exception as ex:
                            _safe_add_paragraph(od, paras[j].text)
                        doc_text_buffer.append(paras[j].text)
                    log_func(f"{e['doc']} 추출 완료")

                extracted_list.append({
                    "doc": e["doc"], "company": e["company"], "link": e["link"],
                    "title": title,
                    "content": "\n".join(doc_text_buffer) if doc_text_buffer else "Conclusion 섹션을 찾지 못했습니다.",
                    "full_content": ("\n".join(full_text_buffer))[:30000] if full_text_buffer else "원문 텍스트를 추출하지 못했습니다."
                })
            except Exception as ex:
                _safe_add_paragraph(od, f"오류 - {e['doc']}: {ex}")
                log_func(str(ex))
                if not extracted_list or extracted_list[-1].get("doc") != e["doc"]:
                    extracted_list.append({
                        "doc": e["doc"], "company": e["company"], "link": e["link"],
                        "title": "(처리 오류)",
                        "content": f"문서 처리 중 오류 발생: {str(ex)[:200]}",
                        "full_content": ""
                    })
            finally:
                # v2.9: 이 문서의 zip/압축폴더를 즉시 해제 (RAM 피크 감소).
                # continue로 빠져나가는 경로에서도 finally는 반드시 실행됨.
                _release_doc_workspace(fp, ed)

            if idx < len(download_results):
                try:
                    od.add_page_break()
                except Exception:
                    pass

        st.session_state.extracted_data = extracted_list
        _build_notebooklm_txt(extracted_list)
        bio = io.BytesIO()
        od.save(bio)
        bio.seek(0)
        return bio
    # _managed_run_dir 컨텍스트 종료 시 temp_dir 자동 정리 (정상/예외 모두)


def _build_notebooklm_txt(extracted_list):
    txt = ["=== 3GPP Contributions Analysis Input ==="]
    for item in extracted_list:
        doc_type = "[CR]" if item.get("is_cr") else "[기고문]"
        txt.append(f"\n\n{'='*50}")
        txt.append(f"문서번호: {item['doc']}  {doc_type}")
        txt.append(f"회사: {item['company']}")
        txt.append(f"제목: {item['title']}")
        txt.append(f"{'='*50}")
        txt.append(item.get('content', ''))
        full = item.get('full_content', '')
        if full and full != item.get('content', ''):
            txt.append(f"\n--- 원문 전체 ---")
            txt.append(full[:5000])
    st.session_state.notebooklm_txt = "\n".join(txt)


# ==========================================
# Output 2 — TF-IDF parse_and_summarize
# ==========================================
class TFIDFEmbedder:
    def __init__(self, max_features=3000, ngram_range=(1, 2)):
        self.v = TfidfVectorizer(
            max_features=max_features, ngram_range=ngram_range,
            lowercase=True, stop_words="english", strip_accents="unicode",
            token_pattern=r"\b[a-zA-Z]{2,}\b",
        )
        self.fitted = False

    def encode(self, texts):
        if isinstance(texts, str): texts = [texts]
        proc = [_RE_SPACES.sub(" ", _RE_NONWORD.sub(" ", t.lower())).strip() for t in texts]
        if not self.fitted:
            self.v.fit(proc)
            self.fitted = True
        return self.v.transform(proc).toarray()


def parse_and_summarize(in_bio, status_elem, log_func):
    d = Document(in_bio)
    props, pcs, cur = [], {}, None

    for el in d.element.body:
        if el.tag.endswith("tbl"):
            tbl = Table(el, d)
            for r in tbl.rows:
                if r.cells[0].text.strip() == "Company":
                    cur = r.cells[1].text.strip()
        elif el.tag.endswith("p"):
            p = Paragraph(el, d)
            txt = p.text.strip()
            is_target = (txt.lower().startswith("proposal") or
                        txt.lower().startswith("summary of change"))
            if is_target:
                buf, cm = [txt], {cur} if cur else set()
                idx2 = d.element.body.index(el) + 1
                while idx2 < len(d.element.body):
                    sib = d.element.body[idx2]
                    if not sib.tag.endswith("p"): break
                    sp = Paragraph(sib, d)
                    st_text = sp.text.rstrip()
                    if not st_text.strip(): break
                    if (st_text.lower().startswith("proposal") or
                        st_text.lower().startswith("summary of change") or
                        st_text.lower().startswith("reason for change")): break
                    buf.append(st_text)
                    if cur: cm.add(cur)
                    idx2 += 1
                bl = "\n".join(buf)
                props.append(bl)
                pcs[bl] = cm.copy()

    r = Document()
    r.add_heading("Proposal Summary", 0)

    if not props:
        _safe_add_paragraph(r, "No proposals found.")
        bio = io.BytesIO()
        r.save(bio)
        bio.seek(0)
        return bio

    status_elem.text("Generating embeddings & Clustering...")
    em = TFIDFEmbedder()
    emb = em.encode(props)

    N = len(props)
    mn, mx = max(2, N // 5), max(3, N // 2)
    best_diff = float("inf")
    best_lbl = None
    for thr in np.linspace(0.2, 0.8, 13):
        try:
            hac = AgglomerativeClustering(
                n_clusters=None, metric="cosine", linkage="average",
                distance_threshold=thr, compute_full_tree=True,
            )
            lbls = hac.fit_predict(emb)
            cnt = len(set(lbls))
            diff = abs(cnt - (mn + mx) / 2)
            if diff < best_diff:
                best_diff = diff
                best_lbl = lbls
        except: pass
    lbls = best_lbl if best_lbl is not None else np.zeros(N, dtype=int)

    clusters = {}
    for i, l in enumerate(lbls):
        clusters.setdefault(l, {"idxs": [], "cm": set()})
        clusters[l]["idxs"].append(i)
        clusters[l]["cm"].update(pcs[props[i]])

    items = []
    for info in clusters.values():
        idxs = info["idxs"]
        subset = emb[idxs]
        cent = np.mean(subset, axis=0, keepdims=True)
        sims = cosine_similarity(cent, subset)[0]
        rep = props[idxs[int(np.argmax(sims))]]
        cm = sorted(info["cm"])
        items.append({"proposal": rep, "companies": cm, "count": len(cm)})

    items.sort(key=lambda x: x["count"], reverse=True)

    status_elem.text("Creating summary...")
    for it in items:
        _safe_add_paragraph(r, it["proposal"])
        _safe_add_paragraph(r, f"Supporting companies ({it['count']}): " + (", ".join(it["companies"]) if it["companies"] else "(none)"))
        _safe_add_paragraph(r, "")

    bio = io.BytesIO()
    r.save(bio)
    bio.seek(0)
    log_func("Summary 생성 완료")
    return bio


# ==========================================
# v2.4: Gemini AI 분석 — 모델 선택 + 503 폴백 + Direct 임계값 상향
# ==========================================
def _build_doc_inventory(extracted_data):
    lines = []
    for item in extracted_data:
        lines.append(f"  - {item['doc']} (회사: {item['company']})")
    return "\n".join(lines)


# ==========================================
# ★ v3.5: API 오류 분류 (상태코드 정확 추출) ★
# ==========================================
# [배경] v3.4는 '403' in err_str 같은 부분매칭을 썼는데,
#   "429 Quota exceeded ... limit 4032 per day" 처럼 숫자 안에 403이
#   들어 있으면 인증 오류로 오분류됨 → 재시도를 안 하고 즉시 중단되어
#   회복 가능한 상황을 놓침. 상태 코드를 정확히 뽑아 판정한다.
_STATUS_CODE_RE = re.compile(r'(?<![0-9])([1-5][0-9]{2})(?![0-9])')


def _extract_status_codes(err_str):
    """오류 문자열에서 HTTP 상태 코드로 보이는 값만 추출.
    앞뒤에 다른 숫자가 붙은 경우(4032 등)는 제외."""
    try:
        return set(int(x) for x in _STATUS_CODE_RE.findall(err_str or ""))
    except Exception:
        return set()


def classify_api_error(err_str):
    """API 오류를 'auth' / 'quota' / 'overload' / 'other' 로 분류.

    분류 기준:
      auth     — 인증/권한 문제. 재시도해도 소용없으므로 즉시 중단.
      quota    — 한도 초과(429). 대기 후 재시도.
      overload — 서버 과부하(503/500/504). 대기 후 재시도 또는 모델 폴백.
      other    — 그 외.
    """
    s = err_str or ""
    low = s.lower()
    codes = _extract_status_codes(s)

    # 1) 명시적 인증 사유 문구가 있으면 코드와 무관하게 auth
    auth_phrases = (
        "unauthenticated", "api key not valid", "api_key_invalid",
        "access_token_type_unsupported", "permission_denied",
        "invalid authentication credentials", "api key expired",
        "invalid api key",
    )
    if any(p in low for p in auth_phrases):
        return "auth"

    # 2) 상태 코드 기반 판정 (한도/과부하를 인증보다 먼저 봐서 오분류 방지)
    if 429 in codes or "resource_exhausted" in low or "quota" in low:
        return "quota"
    if codes & {500, 502, 503, 504} or "overloaded" in low or "unavailable" in low:
        return "overload"
    if codes & {401, 403}:
        return "auth"
    # 400은 키 문제일 수도, 요청 형식 문제일 수도 있음 → 문구로 이미 위에서 걸러짐
    return "other"


def _call_with_retry_and_fallback(api_key, model_obj, prompt, generation_config,
                                    status_elem, max_retries, start_time,
                                    current_model_name, model_choice):
    """
    v2.4: 503 받으면 Flash로 폴백 후 재시도.
    v2.5 수정: 폴백된 모델 객체도 반환 → 호출자가 다음 호출에 사용 가능.
    Returns: (response, final_model_name, final_model_obj) or (None, None, None)
    """
    cur_model = model_obj
    cur_name = current_model_name
    fallback_attempted = False

    for attempt in range(max_retries):
        try:
            res = _safe_gemini_call(api_key, cur_model, prompt, generation_config=generation_config)
            return res, cur_name, cur_model
        except Exception as e:
            err_str = str(e)
            # v3.5: 상태코드를 정확히 추출해 분류 (부분매칭 오탐 제거)
            kind = classify_api_error(err_str)
            is_429 = (kind == "quota")
            is_503 = (kind == "overload")

            # v3.4: 인증 오류는 재시도해도 소용없으므로 즉시 중단하고
            # 사용자에게 원인을 명확히 알린다.
            if kind == "auth":
                append_log(f"인증 오류(재시도 안 함): {err_str[:200]}")
                raise RuntimeError(f"__AUTH_ERROR__{err_str}")

            # v3.0: 503(서버 과부하) 시 다른 Flash 버전으로 폴백.
            # 기존엔 Pro→Flash 폴백이었으나 Pro를 지원하지 않게 되어,
            # 이제는 현재 Flash가 과부하일 때 사용 가능한 다른 Flash로 옮김.
            if is_503 and not fallback_attempted:
                fallback_attempted = True
                valid_models = _get_cached_models(api_key)
                alt_candidates = [
                    m for m in _list_selectable_models(valid_models)
                    if m.split('/')[-1] != cur_name
                ]
                if alt_candidates:
                    # 버전 높은 순으로 정렬해 차선 Flash 선택
                    alt_candidates.sort(key=lambda m: _flash_version_key(m), reverse=True)
                    alt_target = alt_candidates[0]
                    if status_elem:
                        try:
                            status_elem.text(
                                f"⚠️ 서버 과부하 — 다른 Flash로 전환: {alt_target.split('/')[-1]}"
                            )
                        except Exception:
                            pass
                    append_log(f"503 폴백: {cur_name} → {alt_target}")
                    cur_model = genai.GenerativeModel(alt_target)
                    cur_name = alt_target.split('/')[-1]
                    time.sleep(3)
                    continue  # 즉시 재시도 (대기 없음)

            if is_429 or is_503:
                wait = [30, 60, 120, 180, 240][min(attempt, 4)]
                elapsed = int(time.time() - start_time)
                for cd in range(wait, 0, -1):
                    elapsed = int(time.time() - start_time)
                    if status_elem:
                        try:
                            status_elem.text(
                                f"⚠️ API 한도/과부하 대기 {cd}초 (시도 {attempt+1}/{max_retries}, "
                                f"총 {elapsed//60}분 {elapsed%60}초 경과)"
                            )
                        except Exception:
                            pass
                    time.sleep(1)
            else:
                raise

    return None, None, None


def run_gemini_analysis(extracted_data, status_elem, api_key,
                        model_choice="flash_auto", manual_model_name="",
                        call_interval=8):
    """v2.4: 모델 선택, 호출 간격, 503 폴백 추가."""
    if not extracted_data:
        st.warning("⚠️ 분석할 문서가 없습니다.")
        return False

    # v2.9: 진입부의 genai.configure() 제거.
    # 락을 잡았다 즉시 놓는 configure는 효과가 없음 — 실제 호출 시점의
    # _safe_gemini_call()이 락 안에서 configure+generate를 원자적으로 수행.
    _gemini_start_time = time.time()

    doc_inventory = _build_doc_inventory(extracted_data)
    valid_doc_ids = {item['doc'] for item in extracted_data}
    valid_companies = {item['company'] for item in extracted_data if item['company']}

    text_buffer = []
    for item in extracted_data:
        text_buffer.append(
            f"========== 문서 시작: {item['doc']} ==========\n"
            f"회사: {item['company']}\n"
            f"제목: {item.get('title', '')}\n"
            f"내용:\n{item['content']}\n"
            f"========== 문서 끝: {item['doc']} =========="
        )
    full_text = "\n\n".join(text_buffer)

    MAIN_PROMPT = f"""당신은 3GPP 표준화 회의 기고문을 분석하는 전문가입니다.

아래에 다운로드된 기고문 원문이 제공됩니다. 이 원문만을 근거로 분석하세요.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[절대 규칙 — 할루시네이션 금지]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

1. 아래 [허용된 문서 목록]에 있는 문서 번호만 인용할 수 있습니다.
2. 아래 [허용된 회사 목록]에 있는 회사명만 사용할 수 있습니다.
3. 원문에 명시적으로 적혀 있는 내용만 분석하세요.
4. 어떤 회사가 어떤 제안을 지지하는지는, 해당 회사의 기고문에
   해당 제안이 실제로 기술되어 있을 때만 인정됩니다.
5. 확실하지 않으면 포함하지 마세요.

[허용된 문서 목록]
{doc_inventory}

[허용된 회사 목록]
{', '.join(sorted(valid_companies))}

주의: 다음 회사들은 같은 그룹이므로 하나의 회사로 취급하세요:
- ZTE = Sanechips
- Huawei = HiSilicon
- Nokia = Nokia Shanghai Bell

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[분석 구조 — 이슈 → 입장 → 회사]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

이 분석의 목적은 "어떤 쟁점에서 회사들의 의견이 어떻게 갈리는지"를
한눈에 보여주는 것입니다. 따라서 다음 3계층으로 정리하세요.

[1계층] 이슈(쟁점)
  - 여러 회사가 공통으로 다루는 하나의 기술적 논점.
  - 예: "C-DRX 주기와 DL-WUS 모니터링 주기를 동일하게 설정할 것인가"
  - 뭉뚱그리지 마세요. "에너지 효율 관련" 같은 광범위 이슈 금지.
  - 2개 이상 회사가 다룬 이슈만 포함.

[2계층] 입장(option)
  - 하나의 이슈 안에서 회사들이 취하는 서로 다른 견해를 각각 별도로 나눕니다.
  - 입장은 보통 다음 형태로 나타납니다:
      · 서로 다른 해법 제시 (예: "독립 설정" vs "공통 값" vs "네트워크 결정")
      · 어떤 제안에 대한 찬성 / 반대
      · 조건부 찬성(특정 전제 하에서만 동의)
  - 각 입장마다 지지 회사를 빠짐없이 나열하고 회사 수를 셉니다.
  - 한 회사가 한 이슈에서 하나의 입장에만 속하도록 배정하세요.
    (원문에서 판단이 어려우면 그 회사는 해당 입장에 넣지 마세요.)

[3계층] 근거
  - 각 입장을 뒷받침하는 문서번호와 원문 문구.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[입장 분류 규칙]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

1. 찬성/반대가 명시적으로 드러나는 경우에만 찬성·반대로 분류합니다.
   - 찬성 신호: support, agree, we think X is beneficial, propose to adopt 등
   - 반대 신호: do not support, disagree, object, we see no benefit,
                should not be introduced, unnecessary 등
2. 찬반이 아니라 서로 다른 해법을 제시하는 경우에는 찬성/반대 대신
   각 해법을 별도 입장(Option A / Option B / ...)으로 만드세요.
3. "네트워크 구현에 맡기자", "추가 논의 필요", "RAN1 결정에 따름" 같은
   유보적 견해도 하나의 독립된 입장으로 취급하세요.
4. 원문에 근거가 없는 찬반 추정 금지. 애매하면 그 회사를 넣지 마세요.
5. 1개 회사만 단독 주장한 입장이라도, 그 이슈 자체를 2개 이상 회사가
   다루고 있다면 소수 의견으로 표시하세요. (대립 구도 파악에 중요)
6. CR 문서의 "Summary of change"도 제안으로 취급.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[출력 양식]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

### [순위]. [이슈를 구체적으로 요약한 제목]
* **논의 참여 (총 N개사):** 회사명1, 회사명2, ...
* **쟁점 요약:** 이 이슈에서 무엇이 갈리는지 2-3문장.

* **입장별 분류:**

  **[입장 1] (M개사)** 입장을 한 문장으로 요약
  - 지지: 회사명1, 회사명2, 회사명3
  - 근거:
    - [문서번호] (회사명): 핵심 문구 인용
    - [문서번호] (회사명): 핵심 문구 인용

  **[입장 2] (K개사)** 입장을 한 문장으로 요약
  - 지지: 회사명4, 회사명5
  - 근거:
    - [문서번호] (회사명): 핵심 문구 인용

  **[입장 3] (1개사, 소수의견)** 입장을 한 문장으로 요약
  - 지지: 회사명6
  - 근거:
    - [문서번호] (회사명): 핵심 문구 인용

찬반 구도가 뚜렷한 이슈는 입장 이름을 "찬성"/"반대"로 쓰세요:

  **[찬성] (M개사)** 무엇에 찬성하는지
  - 지지: 회사명1, 회사명2
  - 근거:
    - [문서번호] (회사명): 핵심 문구 인용

  **[반대] (K개사)** 무엇에 반대하는지, 반대 이유
  - 지지: 회사명3, 회사명4
  - 근거:
    - [문서번호] (회사명): 핵심 문구 인용

[정렬 규칙]
- 이슈 순위: 논의 참여 회사 수 내림차순.
- 입장 순위: 각 이슈 안에서 지지 회사 수 내림차순.
- 회사 수는 반드시 실제 나열한 회사 개수와 일치해야 합니다.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[기고문 원문 데이터]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

{full_text}
"""

    MAP_PROMPT_TEMPLATE = """당신은 3GPP 기고문 분석 전문가입니다.

[절대 규칙]
1. 아래 문서에 있는 내용만 추출. 없는 내용 지어내지 마세요.
2. 허용된 문서 번호: {doc_list}
3. 각 제안마다 문서 번호와 회사명 함께 기록.
4. 하나의 기고문에 여러 제안이 있으면 각각 별도로.

[입장(stance) 판정 — 2차 통합에서 찬반 진영을 나누는 데 사용됩니다]
각 제안마다 그 회사가 어떤 태도를 취하는지 반드시 표시하세요.
- 찬성: support / agree / propose to adopt / beneficial 등이 명시된 경우
- 반대: do not support / disagree / object / unnecessary /
        no benefit / should not be introduced 등이 명시된 경우
- 제안: 찬반이 아니라 자기 해법을 새로 제시하는 경우
- 유보: 네트워크 구현에 맡기자 / 추가 논의 필요 / 다른 WG 결정에 따름
원문에 근거가 없으면 "불명확"으로 적고, 임의로 찬반을 추정하지 마세요.

[출력 양식]
- 쟁점: [이 제안이 다루는 논점을 한 문장으로]
- 제안: [원문에 가깝게 기술]
- 입장: [찬성 / 반대 / 제안 / 유보 / 불명확]
- 대상: [찬성·반대인 경우, 무엇에 대한 찬반인지]
- 문서: [문서번호]
- 회사: [회사명]
- 원문 근거: [원문 인용]

[기고문 원문]
{batch_text}"""

    REDUCE_PROMPT_TEMPLATE = """당신은 3GPP 기고문 분석 전문가입니다.

[절대 규칙]
1. 허용된 문서 번호: {doc_list}
2. 허용된 회사: {company_list}
3. 1차 추출 결과에 실제로 있는 내용만 사용. 없는 내용 지어내기 금지.

[통합 구조 — 이슈 → 입장 → 회사]
1계층 이슈: 여러 회사가 공통으로 다루는 하나의 기술적 논점.
            광범위 그룹 금지. 2개 이상 회사가 다룬 이슈만 포함.
2계층 입장: 그 이슈 안에서 갈리는 서로 다른 견해를 각각 분리.
            1차 추출의 "입장" 필드를 활용해 찬성/반대/해법별로 나눔.
            유보(네트워크 결정에 맡김 등)도 독립 입장으로 취급.
            한 회사는 한 이슈에서 하나의 입장에만 배정.
3계층 근거: 문서번호 + 원문 문구.

[입장 분류 규칙]
1. 찬반이 명시된 경우에만 찬성/반대로 분류.
2. 서로 다른 해법 제시면 Option 형태의 별도 입장으로.
3. 1개 회사 단독 입장이라도 이슈 자체를 2개 이상 회사가 다루면
   소수의견으로 표시 (대립 구도 파악에 중요).
4. 회사 수는 실제 나열한 회사 개수와 반드시 일치.

[출력 양식]
### [순위]. [이슈를 구체적으로 요약한 제목]
* **논의 참여 (총 N개사):** 회사1, 회사2, ...
* **쟁점 요약:** 무엇이 갈리는지 2-3문장.

* **입장별 분류:**

  **[입장 1] (M개사)** 입장 한 문장 요약
  - 지지: 회사1, 회사2
  - 근거:
    - [문서번호] (회사명): 핵심 문구

  **[입장 2] (K개사)** 입장 한 문장 요약
  - 지지: 회사3
  - 근거:
    - [문서번호] (회사명): 핵심 문구

찬반 구도가 뚜렷하면 입장 이름을 "찬성"/"반대"로 사용하세요.

[정렬]
- 이슈: 논의 참여 회사 수 내림차순
- 입장: 각 이슈 안에서 지지 회사 수 내림차순

[1차 추출 결과]
{intermediate_text}"""

    status_elem.text("🧠 모델을 선택하고 있습니다...")
    try:
        # v2.4: 모델 선택 (매번 재선택, 캐시 안 함)
        target, display_or_err = _resolve_model_for_choice(api_key, model_choice, manual_model_name)
        if not target:
            st.error(
                f"❌ **모델 선택 실패:** {display_or_err}\n\n"
                f"위에서 **'⚙️ 수동 선택'** 모드로 전환하여 사용 가능한 모델을 직접 골라주세요."
            )
            return False

        model_display = display_or_err
        # v3.0 최종 방어선: 어떤 경로로도 유료 모델이 여기 도달하면 중단
        if _is_blocked_model(target):
            st.error(
                f"❌ **차단됨:** '{model_display}'은(는) 유료 결제가 필요한 모델입니다.\n\n"
                f"이 앱은 예기치 않은 과금을 막기 위해 무료 Flash 계열만 사용합니다."
            )
            return False
        model = genai.GenerativeModel(target)
        strict_config = {"temperature": 0.0}

        # v3.2: status_elem.text()는 마크다운을 렌더링하지 않아 **가 그대로
        # 보이므로 markdown()으로 표시. 실제 선택된 모델명을 그대로 노출.
        try:
            status_elem.markdown(f"🧠 **선택된 모델: `{model_display}`** — 분석을 시작합니다...")
        except Exception:
            status_elem.text(f"🧠 선택된 모델: {model_display} — 분석을 시작합니다...")
        append_log(f"선택된 모델: {target}")

        total_docs = len(extracted_data)
        response = None
        final_model = model_display
        doc_list_str = ", ".join(sorted(valid_doc_ids))
        company_list_str = ", ".join(sorted(valid_companies))

        # v2.4: Direct 임계값 20 → 50으로 상향
        DIRECT_THRESHOLD = 50

        if total_docs > DIRECT_THRESHOLD:
            # Map-Reduce
            batch_size = DIRECT_THRESHOLD
            total_batches = (total_docs + batch_size - 1) // batch_size
            intermediate = []
            # v3.3: 부분 실패 추적 + 시간 예산 + 조기 중단
            failed_batches = []
            consecutive_failures = 0
            budget_exceeded = False
            early_stop_reason = ""
            try:
                _budget_sec = int(st.session_state.get(
                    "ai_time_budget_min", DEFAULT_TIME_BUDGET_MIN)) * 60
            except Exception:
                _budget_sec = DEFAULT_TIME_BUDGET_MIN * 60

            for i in range(total_batches):
                # v3.3: 배치 시작 전 시간 예산 검사
                if time.time() - _gemini_start_time > _budget_sec:
                    budget_exceeded = True
                    early_stop_reason = (
                        f"시간 예산({_budget_sec // 60}분) 초과 — "
                        f"{i}/{total_batches} 배치까지만 처리"
                    )
                    append_log(f"⏱️ {early_stop_reason}")
                    for j in range(i, total_batches):
                        failed_batches.append(j + 1)
                    break

                # v3.3: 연속 실패 시 조기 중단 (일일 한도 소진 추정)
                if consecutive_failures >= 2:
                    early_stop_reason = (
                        "연속 2개 배치 실패 — API 일일 한도 소진 가능성. "
                        f"{i}/{total_batches} 배치까지만 처리"
                    )
                    append_log(f"🛑 {early_stop_reason}")
                    for j in range(i, total_batches):
                        failed_batches.append(j + 1)
                    break

                status_elem.text(f"🚀 1차 추출 [{i+1}/{total_batches}]")
                batch = extracted_data[i*batch_size:(i+1)*batch_size]
                bt = "\n\n".join([
                    f"========== {it['doc']} ({it['company']}) ==========\n"
                    f"제목: {it.get('title', '')}\n"
                    f"{it['content']}"
                    for it in batch
                ])
                batch_docs = ", ".join([it['doc'] for it in batch])
                mp = MAP_PROMPT_TEMPLATE.format(doc_list=batch_docs, batch_text=bt)

                res, ret_name, ret_model = _call_with_retry_and_fallback(
                    api_key, model, mp, strict_config,
                    status_elem, 5, _gemini_start_time,
                    final_model, model_choice
                )
                # v2.5: 폴백된 모델 객체로 영구 교체 (다음 배치도 같은 모델 사용)
                if ret_name:
                    final_model = ret_name
                if ret_model is not None and ret_model is not model:
                    model = ret_model
                    append_log(f"모델 영구 교체됨: → {final_model}")

                batch_ok = False
                if res is not None:
                    try:
                        if res.text and len(res.text.strip()) > 10:
                            intermediate.append(res.text)
                            batch_ok = True
                    except (ValueError, AttributeError):
                        append_log(f"배치 {i+1}: 응답 텍스트 접근 실패 (safety filter?)")
                else:
                    append_log(f"배치 {i+1} 실패 (5회 재시도 소진)")

                if batch_ok:
                    consecutive_failures = 0
                else:
                    consecutive_failures += 1
                    failed_batches.append(i + 1)

                # v2.4: 호출 간격 사용자 설정
                if i < total_batches-1:
                    for cd in range(call_interval, 0, -1):
                        status_elem.text(f"⏳ 배치 간 대기 {cd}초 ({i+1}/{total_batches} 완료)")
                        time.sleep(1)

            if not intermediate:
                elapsed = int(time.time() - _gemini_start_time)
                _extra = f"\n\n**중단 사유:** {early_stop_reason}" if early_stop_reason else ""
                st.error(
                    f"❌ **{elapsed//60}분 동안 시도했으나 결과를 받지 못했습니다.**{_extra}\n\n"
                    f"**해결 방법:**\n"
                    f"- 호출 간격을 더 길게 (사이드바 슬라이더)\n"
                    f"- 분석 시간 예산을 늘리기 (사이드바)\n"
                    f"- 새 키 발급 후 재시도\n"
                    f"- 문서 수를 줄여서 재시도\n"
                    f"- NotebookLM(아래 섹션)을 대안으로 사용"
                )
                return False

            # v3.3: 부분 실패를 조용히 넘기지 않고 사용자에게 명시
            if failed_batches:
                _ok = total_batches - len(failed_batches)
                st.warning(
                    f"⚠️ **전체 {total_batches}개 배치 중 {len(failed_batches)}개가 처리되지 않았습니다.** "
                    f"({_ok}개 배치의 문서만 분석에 반영됨)\n\n"
                    + (f"사유: {early_stop_reason}\n\n" if early_stop_reason else "")
                    + "결과에 일부 문서가 누락되어 있으니 참고하세요."
                )
                append_log(
                    f"부분 실패: 배치 {failed_batches} 미반영 "
                    f"({_ok}/{total_batches} 성공)"
                )

            status_elem.text("🧠 최종 병합 분석 중...")
            fi = "\n\n=== 배치 구분 ===\n\n".join(intermediate)
            rp = REDUCE_PROMPT_TEMPLATE.format(
                doc_list=doc_list_str,
                company_list=company_list_str,
                intermediate_text=fi,
            )
            response, ret_name, ret_model = _call_with_retry_and_fallback(
                api_key, model, rp, strict_config,
                status_elem, 5, _gemini_start_time,
                final_model, model_choice
            )
            if ret_name:
                final_model = ret_name
            if ret_model is not None and ret_model is not model:
                model = ret_model
        else:
            # Direct analysis
            response, ret_name, ret_model = _call_with_retry_and_fallback(
                api_key, model, MAIN_PROMPT, strict_config,
                status_elem, 5, _gemini_start_time,
                final_model, model_choice
            )
            if ret_name:
                final_model = ret_name
            if ret_model is not None and ret_model is not model:
                model = ret_model

        status_elem.text("🔍 응답 확인 중...")

        if response is None:
            st.error(
                "❌ **API 응답을 받지 못했습니다.**\n\n"
                "한도 소진 또는 일시적 과부하 상황입니다.\n\n"
                "**해결:** 모델을 Flash로 변경하거나, 새 키 발급 후 재시도."
            )
            return False

        if not hasattr(response, 'text'):
            st.error("❌ **API가 빈 응답을 반환했습니다.**")
            return False

        try:
            result_text = response.text
        except (ValueError, AttributeError) as e:
            st.error(
                "❌ **AI 응답이 안전 필터에 의해 차단되었습니다.**\n\n"
                "**해결:** 다시 시도하거나 NotebookLM 사용."
            )
            append_log(f"Gemini safety filter: {e}")
            return False

        if not result_text or len(result_text.strip()) < 50:
            st.error("❌ **AI 응답이 너무 짧습니다.**")
            return False

        status_elem.text("✅ AI 분석 완료! 결과 문서를 생성하고 있습니다...")

        cited_docs = extract_doc_ids(result_text)
        hallucinated = cited_docs - valid_doc_ids
        if hallucinated:
            result_text += f"\n\n---\n⚠️ **검증 경고:** 다음 문서 번호는 다운로드된 파일 목록에 없습니다 (할루시네이션 가능성): {', '.join(sorted(hallucinated))}"

        # v3.3: 부분 실패가 있었다면 결과 문서 최상단에도 명시
        try:
            if failed_batches:
                _ok = total_batches - len(failed_batches)
                result_text = (
                    f"> ⚠️ **주의:** 전체 {total_batches}개 배치 중 {_ok}개만 분석에 반영되었습니다. "
                    f"일부 문서가 누락된 결과입니다."
                    + (f" (사유: {early_stop_reason})" if early_stop_reason else "")
                    + "\n\n" + result_text
                )
        except NameError:
            pass  # Direct 분석 경로에서는 failed_batches가 없음

        doc = Document()
        doc.add_heading(f"AI 정밀 분석 요약 ({final_model})", 0)
        _safe_add_paragraph(doc, f"분석 대상: {total_docs}개 문서")
        _safe_add_paragraph(doc, f"분석 모델: {final_model} (temperature=0.0)")
        _safe_add_paragraph(doc, "")
        for line in result_text.split('\n'):
            if re.match(r'^(#+)?\s*\d+\.|^###', line.strip()):
                try:
                    p = doc.add_paragraph()
                    cleaned = re.sub(r'^\s*#+\s*', '', line).strip()
                    p.add_run(_xml_safe(cleaned)).bold = True
                except Exception:
                    _safe_add_paragraph(doc, line)
            elif line.strip().startswith('* **'):
                _safe_add_paragraph(doc, line.strip())
            elif line.strip().startswith('- [') or line.strip().startswith('- '):
                _safe_add_paragraph(doc, line.strip())
            else:
                _safe_add_paragraph(doc, line)
        bio = io.BytesIO()
        doc.save(bio)
        st.session_state.ai_summary_bytes = bio.getvalue()
        st.session_state.ai_summary_text = result_text
        st.session_state.ai_model_name = final_model
        st.session_state.ai_summary_generated = True
        status_elem.text("✅ 완료! 아래에서 결과를 확인하세요.")
        st.rerun()
        return True

    except Exception as e:
        err = str(e)
        if GEMINI_API_KEY and GEMINI_API_KEY in err:
            err = err.replace(GEMINI_API_KEY, "***HIDDEN***")
        if api_key and api_key in err:
            err = err.replace(api_key, "***HIDDEN***")
        # v3.5: 분류 함수로 통일 (부분매칭 오탐 제거)
        _is_auth_marked = err.startswith("__AUTH_ERROR__")
        clean = err.replace("__AUTH_ERROR__", "")
        kind = "auth" if _is_auth_marked else classify_api_error(clean)

        if kind == "auth":
            st.error(
                "❌ **API 키 인증에 실패했습니다.** (한도 문제가 아닙니다)\n\n"
                "**확인해 보세요:**\n"
                "1. 키를 전체 복사했는지 (앞뒤 잘림·공백 없이)\n"
                "2. Google AI Studio에서 그 키가 **사용 설정(Enabled)** 상태인지\n"
                "3. 새 형식(`AQ.`) 키를 쓰고 계신다면 — 일부 계정에서 이 형식이 "
                "서버에 거부되는 사례가 보고되고 있습니다. "
                "기존 `AIza` 형식 키가 있다면 그것으로 시도해보세요.\n\n"
                "자세한 오류는 아래 '처리 로그'에서 확인할 수 있습니다."
            )
            append_log(f"인증 실패 상세: {clean[:300]}")
        elif kind == "quota":
            st.error(
                "❌ **API 한도가 초과되었습니다.**\n\n"
                "**해결:** 잠시 후 재시도하거나, 다른 프로젝트에서 발급한 새 키를 사용하세요.\n"
                "(같은 프로젝트의 키를 여러 개 만들어도 한도는 공유됩니다)"
            )
            append_log(f"한도 초과: {clean[:300]}")
        elif kind == "overload":
            st.error(
                "❌ **Gemini 서버가 일시적으로 혼잡합니다.**\n\n"
                "사용자 키 문제가 아니며, 잠시 후 다시 시도하면 대개 해결됩니다."
            )
            append_log(f"서버 과부하: {clean[:300]}")
        else:
            st.error("❌ **API 오류가 발생했습니다.** 자세한 내용은 '처리 로그'를 확인하세요.")
            append_log(f"Gemini error (sanitized): {clean[:300]}")
    return False


# ==========================================
# 심층 분석
# ==========================================

def _parse_ai_summary_into_proposals(ai_summary_text):
    if not ai_summary_text:
        return []
    parts = re.split(r'\n(?=###\s)', ai_summary_text)
    proposals = []
    for part in parts:
        part = part.strip()
        if not part.startswith("###"):
            continue
        lines = part.split('\n', 1)
        header = lines[0].strip()
        body = lines[1].strip() if len(lines) > 1 else ""
        doc_ids = extract_doc_ids(part)
        proposals.append({
            "header": header,
            "body": body,
            "doc_ids": doc_ids,
            "full_block": part,
        })
    return proposals


def _select_docs_for_deep_analysis(doc_ids, extracted_data, max_docs=5):
    matching = [item for item in extracted_data if item.get("doc") in doc_ids]

    def tier_of(company):
        if not company:
            return 3
        if company in MAJOR_VENDORS_TIER1:
            return 1
        if company in MAJOR_VENDORS_TIER2:
            return 2
        return 3

    matching.sort(key=lambda x: tier_of(x.get("company", "")))

    seen_companies = set()
    selected = []
    for item in matching:
        comp = item.get("company", "")
        if comp in seen_companies:
            continue
        seen_companies.add(comp)
        selected.append({
            "doc": item.get("doc", ""),
            "company": comp,
            "full_content": item.get("full_content", ""),
            "tier": tier_of(comp),
        })
        if len(selected) >= max_docs:
            break

    return selected


def run_deep_analysis(proposal_header, proposal_body, selected_docs, api_key,
                       model_choice="flash_auto", manual_model_name=""):
    if not selected_docs:
        return False, "분석 가능한 문서가 없습니다."

    try:
        # v2.9: 진입부 configure 제거 (위와 동일 이유 — 실제 호출 시 재설정됨)

        # v3.0: 무료 Flash 계열만 사용
        target, display_or_err = _resolve_model_for_choice(api_key, model_choice, manual_model_name)
        if not target:
            return False, f"모델 선택 실패: {display_or_err}"
        # v3.0 최종 방어선
        if _is_blocked_model(target):
            return False, f"'{display_or_err}'은(는) 유료 모델이라 사용할 수 없습니다."

        model = genai.GenerativeModel(
            model_name=target,
            generation_config={"temperature": 0.3, "max_output_tokens": 6000},
            safety_settings=[
                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
            ],
        )

        def _smart_truncate(text, max_chars=10000):
            if len(text) <= max_chars:
                return text
            lines = text.split('\n')
            important_lines = []
            other_lines = []
            for line in lines:
                low = line.lower()
                if any(kw in low for kw in ["proposal", "conclusion", "observation",
                                              "summary", "recommendation", "way forward"]):
                    important_lines.append(line)
                else:
                    other_lines.append(line)
            result = "\n".join(important_lines)
            if len(result) > max_chars:
                return result[:max_chars]
            remaining = max_chars - len(result) - 50
            if remaining > 0 and other_lines:
                other_text = "\n".join(other_lines)
                result += "\n\n[기타 본문 발췌]\n" + other_text[:remaining]
            return result

        doc_texts = []
        for d in selected_docs:
            content = _smart_truncate(d["full_content"], 10000)
            doc_texts.append(f"\n━━━ 문서: {d['doc']} ({d['company']}) ━━━\n{content}")

        combined_content = "\n".join(doc_texts)

        prompt = f"""당신은 3GPP 표준화 회의 기고문을 심층 분석하는 전문가입니다.

아래 제안에 대해, 지지 회사들의 원문을 바탕으로 **근거와 반박 논리**를 분석해주세요.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[분석 대상 제안]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{proposal_header}

{proposal_body}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[지지 회사 원문 ({len(selected_docs)}개 기고문)]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{combined_content}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[출력 규칙]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
1. [근거] 섹션: 오직 위에 제공된 원문 내용만 사용.
2. [진영 구도] 섹션: 원문에 나타난 입장만 사용. 추정 금지.
3. [반박] 섹션: **반드시 작성** (절대 생략 불가). "⚠️ 추론" 태그 필수.
4. [전략적 함의] 섹션: **반드시 작성**.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
[출력 양식]
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

### 🗳️ 진영 구도

원문에서 각 회사가 취한 입장을 분류하세요. 회사 수를 반드시 표시합니다.
찬반이 아니라 서로 다른 해법이면 Option 형태로 나누세요.
원문에 입장이 명확하지 않은 회사는 "입장 불명확"에 넣으세요.

- **[찬성 / 입장 A] (N개사):** 회사1, 회사2, ...
  - 요지: 한 문장
- **[반대 / 입장 B] (M개사):** 회사3, 회사4, ...
  - 요지: 한 문장
- **[유보·조건부] (K개사):** 회사5
  - 요지: 한 문장
- **[입장 불명확] (L개사):** 회사6

### 🧠 주장의 근거 및 논리

**핵심 논거:**
1. [논거 1 — 원문 기반]
2. [논거 2]
3. [논거 3]

**회사별 논리 전개:**
- **[회사A] ([문서번호]):** 논리 전개 요약
- **[회사B] ([문서번호]):** ...

### ⚡ 가능한 반박 논리 (필수 섹션)

**(a) 원문에 나타난 반대 의견:**
- [있으면 인용. 없으면 "원문에는 명시적 반대 없음" 기재 후 (b)로]

**(b) ⚠️ 추론 — 기술적/구현적 관점의 반박:**
1. **[관점 1 — 예: 구현 복잡도]**: 우려사항 상세
2. **[관점 2 — 예: 오버헤드]**: ...
3. **[관점 3 — 예: 호환성]**: ...

### 🎯 전략적 함의 (필수 섹션)
- **채택 시 영향:** [어떤 회사군에 유리한지]
- **거절 시 영향:** [대체 방안]
- **핵심 관전 포인트:** [주의 깊게 볼 논점]
"""

        response = _safe_gemini_call(api_key, model, prompt)
        if not response or not hasattr(response, 'text'):
            return False, "API 응답을 받지 못했습니다."

        try:
            result = response.text
        except (ValueError, AttributeError):
            return False, "AI 응답이 안전 필터에 의해 차단되었습니다."

        if not result or len(result.strip()) < 50:
            return False, "AI 응답이 너무 짧습니다."

        # v3.2: 어떤 모델로 분석했는지 결과 하단에 명시
        result = f"{result}\n\n---\n*분석 모델: `{display_or_err}`*"
        return True, result

    except Exception as e:
        err = str(e)
        if GEMINI_API_KEY and GEMINI_API_KEY in err:
            err = err.replace(GEMINI_API_KEY, "***HIDDEN***")
        if api_key and api_key in err:
            err = err.replace(api_key, "***HIDDEN***")
        # v3.5: 메인 분석과 동일한 분류 기준 적용
        _marked = err.startswith("__AUTH_ERROR__")
        clean = err.replace("__AUTH_ERROR__", "")
        kind = "auth" if _marked else classify_api_error(clean)
        if kind == "auth":
            return False, ("API 키 인증 실패. 키가 유효한지 확인해주세요. "
                           "(한도 문제가 아닙니다)")
        if kind == "quota":
            return False, "API 한도 초과. 잠시 후 재시도하거나 다른 프로젝트의 키를 사용하세요."
        if kind == "overload":
            return False, "Gemini 서버가 일시적으로 혼잡합니다. 잠시 후 재시도해주세요."
        return False, f"오류: {clean[:200]}"


# ==========================================
# Streamlit UI
# ==========================================
st.sidebar.title("📡 3GPP Analyzer v2")
st.sidebar.caption("기본 분석 + Gemini AI 강화")
st.sidebar.markdown("---")

# v2.4: 사이드바에 호출 간격 슬라이더
# v2.5 수정: key 파라미터로 session_state 자동 관리. 강제 할당 패턴 제거.
with st.sidebar.expander("⚙️ Gemini API 고급 설정", expanded=False):
    st.slider(
        "배치 사이 대기 시간 (초)",
        min_value=5, max_value=30,
        key="ai_call_interval",  # session_state["ai_call_interval"]에 자동 저장
        help="짧게 = 빠르지만 한도 초과 위험. 길게 = 안전하지만 느림. Flash는 8초, Pro는 15초 권장."
    )
    st.caption(f"현재: {st.session_state.ai_call_interval}초")

    # v3.3: 분석 시간 예산
    st.slider(
        "분석 시간 예산 (분)",
        min_value=5, max_value=50,
        key="ai_time_budget_min",
        help="이 시간을 넘기면 남은 배치를 중단하고 그때까지 모은 결과로 분석을 마칩니다."
    )
    st.caption(
        f"현재: {st.session_state.ai_time_budget_min}분 — "
        f"Cloud Run 등 배포 환경의 요청 타임아웃보다 짧게 설정하세요."
    )

page = st.sidebar.radio("메뉴", ["🚀 통합 분석기", "⚙️ 설정", "ℹ️ 가이드"])

if page == "⚙️ 설정":
    st.title("⚙️ 서버 설정")
    st.subheader("🔒 보안 안내")
    st.info(
        "API 키와 Cloud Function URL은 서버 환경변수 또는 Streamlit Secrets에 저장되며, "
        "사용자 브라우저에 노출되지 않습니다."
    )
    st.subheader("Gemini API Key")
    st.info(f"상태: {'✅ 설정됨 (서버에 안전하게 저장)' if GEMINI_API_KEY else '❌ 미설정'}")
    if not GEMINI_API_KEY:
        st.code('# .streamlit/secrets.toml\nGEMINI_API_KEY = "AQ.Ab..."', language="toml")
    st.subheader("Cloud Function URL")
    if CLOUD_FUNCTION_URL:
        masked_url = CLOUD_FUNCTION_URL[:40] + "..." if len(CLOUD_FUNCTION_URL) > 40 else CLOUD_FUNCTION_URL
        st.info(f"상태: ✅ 설정됨 ({masked_url})")
    else:
        st.info("상태: ⚠️ 미설정")

elif page == "ℹ️ 가이드":
    st.title("ℹ️ 사용 가이드")
    st.header("🔰 기본 사용법")
    st.markdown("""
**1단계:** 🔍 회의 번호로 자동 조회 → Working Group과 회의 번호 입력
**2단계:** 📋 Agenda 선택 → 🚀 기본 분석 실행
**3단계:** ✨ Gemini AI 정밀 분석 (선택)
    """)

    st.markdown("---")
    st.header("🔑 Gemini API 키 발급 가이드")
    st.markdown("""
**1단계:** [Google AI Studio - API 키 발급 페이지](https://aistudio.google.com/app/apikey)
**2단계:** **'Create API key'** → **'Create API key in new project'** ⚠️ 반드시 in new project 선택!
**3단계:** 발급된 키를 복사하여 분석기에 붙여넣기 (`AQ.Ab...` 또는 `AIzaSy...`)

✅ **완전 무료**, 카드 등록 불필요
    """)

    st.markdown("---")
    st.header("📊 모델 선택 가이드 (v3.5)")
    st.markdown("""
**🟢 Flash 자동 (권장):**
- 분당 10회, 일 250회까지 무료
- 빠르고 안정적
- 대부분의 분석에 충분

**🟡 Pro 자동:**
- 분당 5회, 일 100회 (한도 빠르게 소진)
- 더 똑똑한 추론
- 503 ServiceUnavailable 자주 발생 → 자동으로 Flash 폴백

**⚙️ 수동 선택:**
- 가용 모델 목록에서 직접 선택
- 특정 버전 강제 사용 가능
    """)

elif page == "🚀 통합 분석기":
    st.title("🚀 3GPP 기고문 통합 분석기")
    st.caption("Output 1·2는 기본 | Output 3 Gemini는 선택")

    st.caption(
        "v3.5 — 새 형식(AQ.) API 키 지원  \n"
        "· 쟁점별로 찬성·반대 회사를 나눠서 표시  \n"
        "· 인증 실패·한도 초과·서버 혼잡을 구분해 안내"
    )

    # Step 1: Input
    st.header("1️⃣ 데이터 입력")
    if CLOUD_FUNCTION_URL:
        st.success("☁️ Cloud Function 연결됨 — 클라우드에서 처리합니다.")

    input_method = st.radio(
        "입력 방식:",
        ("🔍 회의 번호로 자동 조회", "Excel 파일 업로드", "링크 직접 입력"),
        horizontal=True,
    )
    entries = []

    if input_method == "🔍 회의 번호로 자동 조회":
        col_wg, col_num = st.columns([1, 2])
        with col_wg:
            wg = st.selectbox("Working Group:", list(WG_FTP_MAP.keys()))

        if st.session_state.get("_last_selected_wg") != wg:
            st.session_state["_last_selected_wg"] = wg
            st.session_state["resolved_folder"] = None
            st.session_state["agenda_dict"] = {}
            st.session_state["all_entries"] = []

        with col_num:
            meeting_num_input = st.text_input(
                "회의 번호 입력 후 Enter ↵ (예: 133bis, 122, 168):",
                placeholder="133bis",
            )

        if meeting_num_input and meeting_num_input.strip():
            meeting_num = meeting_num_input.strip()
            st.info("✅ 회의 번호 입력 완료 — 아래 **Agenda 불러오기** 버튼을 클릭하세요.")

            if st.button("📋 Agenda 불러오기", type="primary"):
                with st.spinner(f"{wg}#{meeting_num} 폴더 검색 중..."):
                    meeting_folder = resolve_meeting_folder(wg, meeting_num)
                    if meeting_folder:
                        st.session_state["resolved_folder"] = meeting_folder
                        agenda_dict, all_entries = fetch_tdoc_list_xlsx(wg, meeting_folder)
                        st.session_state.agenda_dict = agenda_dict
                        st.session_state.all_entries = all_entries
                        if not agenda_dict:
                            st.error(f"❌ TDoc 리스트를 찾지 못했습니다.")
                    else:
                        st.error(f"❌ {wg}#{meeting_num} 폴더를 찾지 못했습니다.")
                        st.session_state.agenda_dict = {}
                        st.session_state.all_entries = []

            if st.session_state.get("resolved_folder"):
                st.caption(f"📂 `ftp/{WG_FTP_MAP.get(wg, '')}/{st.session_state.get('resolved_folder', '')}/Docs/`")

        if st.session_state.agenda_dict:
            agenda_items = sorted(st.session_state.agenda_dict.keys())
            st.success(f"✅ {len(agenda_items)}개 agenda, 총 {len(st.session_state.all_entries)}개 문서 발견")

            st.markdown("#### 👇 분석할 Agenda를 선택하세요")
            selected_agenda = st.selectbox(
                "Agenda 선택:",
                agenda_items,
                format_func=lambda x: f"{x} ({len(st.session_state.agenda_dict[x])}개 문서)",
                label_visibility="collapsed",
            )

            if selected_agenda:
                entries = st.session_state.agenda_dict[selected_agenda]
                st.session_state["selected_agenda_name"] = selected_agenda
                st.info(f"📄 **{selected_agenda}** — {len(entries)}개 문서가 분석 대상입니다.")

                with st.expander(f"문서 목록 미리보기 ({len(entries)}개)", expanded=False):
                    for e in entries[:30]:
                        st.text(f"  {e['doc']}  |  {e['company']}")
                    if len(entries) > 30:
                        st.caption(f"  ... 외 {len(entries)-30}개")

    elif input_method == "Excel 파일 업로드":
        uploaded = st.file_uploader("엑셀(.xlsx) — 1열: 문서번호(하이퍼링크), 3열: 회사명", type=["xlsx","xls"])
        if uploaded:
            entries, skipped = read_excel_from_bytes(uploaded)
            st.info(f"총 {len(entries)}개 문서 인식")
            # v3.3: 제외된 행을 사용자에게 명시 (조용히 누락되지 않도록)
            if skipped:
                with st.expander(f"⚠️ 제외된 행 {len(skipped)}개 — 사유 확인", expanded=False):
                    for _d, _why in skipped[:50]:
                        st.text(f"  {_d}  —  {_why}")
                    if len(skipped) > 50:
                        st.caption(f"  ... 외 {len(skipped)-50}개")
                    st.caption(
                        "하이퍼링크가 없는 행은 문서 경로를 확정할 수 없어 제외됩니다. "
                        "'🔍 회의 번호로 자동 조회' 방식을 이용하시면 자동으로 경로가 채워집니다."
                    )
    else:
        raw = st.text_area("3GPP .zip 링크를 한 줄에 하나씩:", height=120)
        if raw:
            rejected = []
            for line in raw.strip().split("\n"):
                url = line.strip()
                if not url:
                    continue
                # v3.3: 3gpp.org 외 도메인 차단
                if not _is_allowed_doc_url(url):
                    rejected.append(url)
                    continue
                docid = url.split("/")[-1].replace(".zip","")
                entries.append({"doc": docid, "company": "Unknown", "link": url})
            st.info(f"총 {len(entries)}개 문서 인식")
            if rejected:
                st.error(
                    f"❌ 허용되지 않은 도메인 {len(rejected)}개 제외 — "
                    f"3gpp.org 도메인만 다운로드할 수 있습니다."
                )
                with st.expander("제외된 링크 보기", expanded=False):
                    for _u in rejected[:20]:
                        st.text(f"  {_u[:100]}")

    # Step 2: 기본 분석
    st.markdown("---")
    st.header("2️⃣ 기본 분석 (Output 1 + 2)")
    st.write("결론(Conclusions) 추출 + TF-IDF 기반 Proposal 요약을 생성합니다.")

    if st.button("🚀 기본 분석 실행 (Run)", type="primary", use_container_width=True):
        if not entries:
            st.warning("먼저 데이터를 입력해주세요.")
        else:
            st.session_state.log_text = ""
            st.session_state.process_done = False
            st.session_state.ai_summary_generated = False
            st.session_state.ai_summary_bytes = None
            st.session_state.ai_summary_text = ""
            st.session_state.out1_bytes = None
            st.session_state.out2_bytes = None
            st.session_state.extracted_data = []
            st.session_state.notebooklm_txt = None
            st.session_state.deep_analysis_cache = {}
            st.session_state.deep_analysis_inflight = set()
            # v2.6 Fix E: PDF skip 카운터 리셋
            st.session_state.pdf_skip_count = 0

            progress_container = st.container()
            with progress_container:
                st.subheader("📊 처리 진행 상황")
                progress_bar = st.progress(0)
                status_text = st.empty()
                step_detail = st.empty()

                status_text.markdown("**📥 Phase 1/2:** 기고문 다운로드 및 결론 추출")
                step_detail.caption(f"총 {len(entries)}개 문서를 3GPP 서버에서 다운로드합니다...")
                out1_bio = extract_all_conclusions(entries, step_detail, progress_bar, append_log)

                status_text.markdown("**🔬 Phase 2/2:** TF-IDF 기반 제안 클러스터링")
                step_detail.caption("단어 빈도 분석으로 유사한 제안을 자동 그룹핑합니다...")
                out2_bio = parse_and_summarize(out1_bio, step_detail, append_log)

                progress_bar.progress(1.0)
                status_text.markdown("**✅ 기본 분석 완료!**")
                step_detail.empty()

            st.session_state.out1_bytes = out1_bio.getvalue()
            st.session_state.out2_bytes = out2_bio.getvalue()
            st.session_state.process_done = True

    if st.session_state.process_done:
        # v2.6 Fix E: PDF 스킵 발생 시 사용자에게 명시적 안내
        # (st.toast 대신 st.warning — 기존 코드 UI 일관성)
        pdf_skipped = st.session_state.get("pdf_skip_count", 0)
        if pdf_skipped > 0:
            st.warning(
                f"⚠️ **{pdf_skipped}개의 PDF 문서가 처리되지 않았습니다** "
                f"(PyMuPDF 라이브러리 미설치).\n\n"
                f"해당 문서는 Output 3 (AI 분석)에서 자동 제외됩니다. "
                f"서버 관리자가 `pip install PyMuPDF`를 실행하면 다음 분석부터 처리됩니다."
            )

        st.success("🎉 기본 분석 완료! Output 1·2를 다운로드하세요.")

        agenda_tag = _safe_filename(st.session_state.get("selected_agenda_name", ""), 30)
        if not agenda_tag:
            agenda_tag = "manual"

        col1, col2 = st.columns(2)
        with col1:
            if st.session_state.out1_bytes:
                st.download_button("📥 Output 1 (Conclusions 취합.docx)",
                    data=st.session_state.out1_bytes,
                    file_name=f"output1_conclusions_{agenda_tag}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True)
        with col2:
            if st.session_state.out2_bytes:
                st.download_button("📥 Output 2 (TF-IDF 요약.docx)",
                    data=st.session_state.out2_bytes,
                    file_name=f"output2_summary_tfidf_{agenda_tag}.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    use_container_width=True)

        # Step 3: Gemini AI 정밀 분석
        st.markdown("---")
        st.header("3️⃣ AI 정밀 분석 (Gemini)")
        st.write("추출된 결론을 Gemini AI로 의미 분석합니다.")

        st.warning(
            "⏱️ **소요 시간 안내:** 무료 Gemini API는 분당 처리량이 제한되어 있습니다. "
            "문서 수에 따라 **약 5~15분** 이상 소요될 수 있습니다."
        )

        # v3.0: 모델 선택 UI — 무료 Flash 계열 전용
        st.markdown("#### 🎯 분석 모델 선택")
        st.info(
            "💚 이 앱은 **무료 티어 Flash 계열만** 사용합니다. "
            "유료 결제가 필요한 Pro 계열은 지원하지 않으므로 토큰 요금이 청구되지 않습니다."
        )
        model_choice_label = st.radio(
            "모델 선택:",
            list(MODEL_DISPLAY_OPTIONS.values()),
            index=0,  # 최신 Flash 자동이 기본
            horizontal=False,
        )
        # 라벨에서 키로 역변환
        model_choice = next(k for k, v in MODEL_DISPLAY_OPTIONS.items() if v == model_choice_label)
        st.session_state.ai_model_choice = model_choice

        manual_model_name = ""
        if model_choice == "manual":
            # 키가 있어야 가용 모델 조회 가능 — 그 전엔 안내만
            st.info("💡 API 키를 아래에 입력한 후 가용 모델 목록을 보여드립니다.")

        api_key_to_use = None

        if GEMINI_API_KEY:
            key_mode = st.radio(
                "API 키 선택:",
                ("🔑 서버 기본 키 사용", "🔐 내 개인 Gemini API 키 사용"),
                horizontal=True,
            )
            if "개인" in key_mode:
                with st.expander("📖 개인 API 키 발급 방법", expanded=True):
                    st.markdown("**1단계:** [Google AI Studio](https://aistudio.google.com/app/apikey)")
                    st.markdown("**2단계:** **Create API key** → **Create API key in new project** ⚠️ 반드시 in new project!")
                    st.markdown("**3단계:** 발급된 키 복사 → 아래 붙여넣기 (`AQ.Ab...` 형식으로 나옵니다)")
                    st.markdown("💡 **결제(billing)는 등록하지 마세요.** 무료 티어 그대로 쓰면 요금이 발생하지 않습니다.")

                personal_key = st.text_input(
                    "개인 Gemini API Key 입력:",
                    type="password",
                    placeholder="AQ.Ab... 또는 AIzaSy...",
                )
                if personal_key and personal_key.strip():
                    cleaned_key, key_err = validate_api_key_format(personal_key)
                    if key_err:
                        st.error(f"❌ {key_err}")
                    else:
                        api_key_to_use = cleaned_key
                else:
                    st.caption("⬆️ 위에 개인 API 키를 입력하세요.")
            else:
                api_key_to_use = GEMINI_API_KEY
        else:
            st.info("서버에 기본 API 키가 설정되어 있지 않습니다.")
            with st.expander("📖 API 키 발급 방법", expanded=True):
                st.markdown("""
**1단계:** [Google AI Studio](https://aistudio.google.com/app/apikey)
**2단계:** **Create API key** → **Create API key in new project**
**3단계:** 발급된 키 복사 (`AQ.Ab...` 또는 `AIzaSy...`)
                """)
            personal_key = st.text_input(
                "Gemini API Key 입력:",
                type="password",
                placeholder="AQ.Ab... 또는 AIzaSy...",
            )
            if personal_key and personal_key.strip():
                cleaned_key, key_err = validate_api_key_format(personal_key)
                if key_err:
                    st.error(f"❌ {key_err}")
                else:
                    api_key_to_use = cleaned_key

        # v3.0: 수동 모델 선택 UI — 무료 Flash 계열만 노출 (유료 모델 원천 차단)
        if api_key_to_use and model_choice == "manual":
            try:
                with st.spinner("가용 모델 조회 중..."):
                    valid_models = _get_cached_models(api_key_to_use)
                selectable = _list_selectable_models(valid_models)
                # 최신 버전이 위로 오도록 정렬
                selectable.sort(key=lambda m: _flash_version_key(m), reverse=True)
                if selectable:
                    short_names = [m.split("/")[-1] for m in selectable]
                    selected_short = st.selectbox(
                        "사용할 모델 선택 (무료 Flash 계열):",
                        short_names,
                        help="유료 결제가 필요한 Pro 계열은 목록에서 제외됩니다."
                    )
                    manual_model_name = selected_short
                    st.session_state.ai_manual_model_name = manual_model_name
                elif valid_models:
                    st.error("이 키로 사용 가능한 무료 Flash 모델이 없습니다.")
                else:
                    st.error("가용 모델 목록을 가져오지 못했습니다.")
            except Exception as e:
                err = str(e)
                if api_key_to_use in err:
                    err = err.replace(api_key_to_use, "***HIDDEN***")
                st.error(f"모델 조회 오류: {err[:200]}")

        if api_key_to_use:
            # 모델 검증 (수동 모드 외에는 자동 매칭 결과 미리보기)
            if model_choice != "manual":
                try:
                    preview_target, preview_display = _resolve_model_for_choice(
                        api_key_to_use, model_choice, ""
                    )
                    if preview_target:
                        st.success(f"✅ 선택된 모델: **{preview_display}**")
                        # v3.2: 왜 이 모델이 뽑혔는지 확인할 수 있도록 후보 목록 제공
                        try:
                            _cands = _list_selectable_models(_get_cached_models(api_key_to_use))
                            _cands = sorted(_cands, key=_flash_version_key, reverse=True)
                            if len(_cands) > 1:
                                with st.expander(
                                    f"🔍 사용 가능한 Flash 모델 {len(_cands)}개 — 선택 근거 확인",
                                    expanded=False
                                ):
                                    st.caption(
                                        "버전이 가장 높은 모델을 자동 선택합니다. "
                                        "(같은 버전이면 Lite보다 일반 Flash 우선)"
                                    )
                                    for _m in _cands:
                                        _n = _model_short_name(_m)
                                        _mark = "**← 선택됨**" if _m == preview_target else ""
                                        st.markdown(f"- `{_n}` {_mark}")
                        except Exception:
                            pass
                    else:
                        st.warning(f"⚠️ {preview_display}")
                except Exception as _pe:
                    # v3.4: 키 입력 직후 인증 실패를 즉시 알려준다.
                    # v3.5: 분류 함수로 통일
                    _pes = str(_pe)
                    if api_key_to_use and api_key_to_use in _pes:
                        _pes = _pes.replace(api_key_to_use, "***HIDDEN***")
                    if classify_api_error(_pes) == "auth":
                        st.error(
                            "❌ **이 키로 인증할 수 없습니다.**\n\n"
                            "- 키를 전체 복사했는지 확인해주세요 (앞뒤 잘림·공백 없이)\n"
                            "- Google AI Studio에서 해당 키가 활성 상태인지 확인해주세요\n"
                            "- 새 형식(`AQ.`) 키는 일부 계정에서 서버가 거부하는 사례가 "
                            "보고되고 있습니다. 기존 `AIza` 키가 있다면 그것으로 시도해보세요."
                        )
                        with st.expander("오류 상세", expanded=False):
                            st.code(_pes[:500])
                    else:
                        st.warning("⚠️ 모델 목록을 확인하지 못했습니다. 분석 시 다시 시도합니다.")

            st.markdown("#### 👇 준비가 되었으면 아래 버튼을 클릭하세요")
            if st.button("✨ Gemini AI 정밀 분석 시작", use_container_width=True, type="primary"):
                gemini_container = st.container()
                with gemini_container:
                    st.subheader("🧠 Gemini AI 분석 진행 상황")
                    gemini_status = st.empty()
                    gemini_detail = st.empty()

                    total_docs = len(st.session_state.extracted_data)
                    if total_docs == 0:
                        st.warning("⚠️ 추출된 문서가 없습니다.")
                    elif total_docs > 50:
                        total_batches = (total_docs + 49) // 50
                        gemini_detail.caption(
                            f"📋 {total_docs}개 문서를 {total_batches}개 배치로 분석. 약 {max(5, total_batches * 2)}분 소요."
                        )
                    else:
                        gemini_detail.caption(
                            f"📋 {total_docs}개 문서를 일괄 분석 (Direct 모드). 약 3~10분 소요."
                        )

                    if total_docs > 0:
                        run_gemini_analysis(
                            st.session_state.extracted_data,
                            gemini_status,
                            api_key_to_use,
                            model_choice=model_choice,
                            manual_model_name=manual_model_name,
                            call_interval=st.session_state.ai_call_interval,
                        )

        if st.session_state.ai_summary_generated:
            st.success("✅ AI 정밀 요약 완료!")
            st.download_button(
                f"📥 Output 3 (AI 정밀 요약 - {st.session_state.ai_model_name}.docx)",
                data=st.session_state.ai_summary_bytes,
                file_name=f"Output3_AI_Summary_{agenda_tag}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                type="primary", use_container_width=True)

            proposals = _parse_ai_summary_into_proposals(st.session_state.ai_summary_text)

            st.markdown("""
<style>
div.stButton > button[kind="primary"] {
    background-color: #FF4B4B;
    color: white;
    border: 2px solid #D32F2F;
    font-weight: bold;
}
</style>
            """, unsafe_allow_html=True)

            with st.expander("👀 AI 분석 결과 미리보기 + 심층 분석", expanded=True):
                if not proposals:
                    st.markdown(st.session_state.ai_summary_text)
                else:
                    st.info(f"💡 총 {len(proposals)}개 쟁점. **'근거 및 반박 논리 분석'** 버튼을 누르면 진영별 심층 분석.")

                    for p_idx, prop in enumerate(proposals):
                        st.markdown(prop["full_block"])

                        if not prop["doc_ids"]:
                            st.caption("ℹ️ 문서 번호가 없어 심층 분석을 할 수 없습니다.")
                            st.markdown("---")
                            continue

                        cache_key = (prop['header'], tuple(sorted(prop['doc_ids'])))
                        # v2.4: 진행 중 키 추적 (더블클릭 방지)
                        is_inflight = cache_key in st.session_state.deep_analysis_inflight
                        is_cached = cache_key in st.session_state.deep_analysis_cache

                        col_btn, col_info = st.columns([2, 3])
                        with col_btn:
                            btn_label = "🔄 분석 중..." if is_inflight else "🔍 근거 및 반박 논리 분석"
                            btn_clicked = st.button(
                                btn_label,
                                key=f"deep_btn_{p_idx}",
                                type="primary",
                                use_container_width=True,
                                disabled=is_inflight,  # v2.4: 진행 중 비활성화
                            )
                        with col_info:
                            preview_docs = _select_docs_for_deep_analysis(
                                prop["doc_ids"], st.session_state.extracted_data, max_docs=5
                            )
                            if preview_docs:
                                preview_str = ", ".join([f"{d['company']}" for d in preview_docs])
                                st.caption(f"📋 분석 대상 {len(preview_docs)}개사: {preview_str}")

                        if btn_clicked and not is_inflight:
                            # v2.4: 진행 중 표시
                            st.session_state.deep_analysis_inflight.add(cache_key)
                            # v2.5 수정: success 변수를 try 진입 전에 초기화
                            # (어느 분기에서 함수 호출에 실패해도 finally에서 안전하게 참조 가능)
                            success = False
                            try:
                                with st.spinner(f"🔍 '{prop['header'][:50]}...' 심층 분석 중..."):
                                    selected_docs = _select_docs_for_deep_analysis(
                                        prop["doc_ids"], st.session_state.extracted_data, max_docs=5
                                    )
                                    if not selected_docs:
                                        st.error("❌ 분석할 문서가 메모리에 없습니다.")
                                    else:
                                        deep_api_key = api_key_to_use or GEMINI_API_KEY
                                        if not deep_api_key:
                                            st.error("❌ API 키가 설정되지 않았습니다.")
                                        else:
                                            success, result = run_deep_analysis(
                                                prop["header"], prop["body"],
                                                selected_docs, deep_api_key,
                                                model_choice=model_choice,
                                                manual_model_name=manual_model_name,
                                            )
                                            if success:
                                                if cache_key in st.session_state.deep_analysis_cache:
                                                    del st.session_state.deep_analysis_cache[cache_key]
                                                elif len(st.session_state.deep_analysis_cache) >= 30:
                                                    oldest = next(iter(st.session_state.deep_analysis_cache))
                                                    del st.session_state.deep_analysis_cache[oldest]
                                                st.session_state.deep_analysis_cache[cache_key] = result
                                            else:
                                                st.error(f"❌ 심층 분석 실패: {result}")
                            finally:
                                # v2.6 Fix I (Supplement 2 통합):
                                # CRITICAL — rerun 전에 inflight 해제. 다음 렌더 사이클에서
                                # 버튼이 "분석 중" 상태로 잘못 disabled되는 것 방지.
                                # finally는 항상 실행되므로 예외 발생해도 안전하게 해제됨.
                                st.session_state.deep_analysis_inflight.discard(cache_key)
                            # v2.6 Fix I: success/실패 무관하게 rerun.
                            # 성공 → 캐시된 결과 표시.
                            # 실패 → 에러 메시지가 이미 표시됐고, 버튼이 다시 enable됨을 즉시 반영.
                            # 무한 루프 위험 없음 — rerun 후 다음 렌더에서는
                            # btn_clicked가 False로 리셋되므로 이 분기 재진입 안 함.
                            if btn_clicked:
                                st.rerun()

                        if is_cached:
                            with st.container():
                                st.markdown("---")
                                st.markdown(f"#### 🔬 심층 분석 결과 — {prop['header'].replace('###','').strip()}")
                                st.markdown(st.session_state.deep_analysis_cache[cache_key])

                        st.markdown("---")

        # Step 4: NotebookLM
        st.markdown("---")
        st.header("4️⃣ Google NotebookLM 활용하기 (대안)")
        st.success("💡 **환각 제로!** NotebookLM은 오직 업로드한 문서 기반으로만 답변.")

        col_a, col_b = st.columns([2, 1])
        with col_a:
            st.markdown("""
**[NotebookLM 장점]**
* **완전 무료**, 토큰 초과 없음
* **초대용량 지원** (노트북당 50개 파일, 파일당 50만 단어)
* **출처 표기 (Citation)**
            """)
        with col_b:
            if st.session_state.notebooklm_txt:
                st.download_button(
                    label="📝 NotebookLM 전용 텍스트(.txt) 다운로드",
                    data=st.session_state.notebooklm_txt.encode('utf-8'),
                    file_name=f"NotebookLM_Conclusions_{agenda_tag}.txt",
                    mime="text/plain",
                    type="primary",
                    use_container_width=True,
                )

        st.markdown("---")
        st.markdown("#### 📋 NotebookLM 사용법")
        st.markdown("1. 위 버튼으로 **.txt** 저장")
        st.markdown("2. [Google NotebookLM](https://notebooklm.google.com/) 접속")
        st.markdown("3. **새 노트북** → .txt 업로드")
        st.markdown("4. 채팅창에 아래 프롬프트 붙여넣기:")

        notebooklm_prompt = """당신은 3GPP 표준화 회의의 전문 기술 분석가입니다.

[분석 지침]
1. 2개 이상 회사가 공통 주장한 제안만 추출
2. CR 문서의 "Summary of change"도 포함
3. 그룹화: 기술적 핵심 의미가 동일하면 묶되, 광범위 그룹 금지
4. 같은 주제라도 방향 다르면 별도 그룹
5. 지지 회사 수 내림차순
6. 회사 그룹: ZTE=Sanechips, Huawei=HiSilicon, Nokia=Nokia Shanghai Bell

[출력]
### [순위]. [제안 제목]
* **지지 회사 (총 N개사):** ...
* **상세 내용:** 2-3문장
* **근거 문서:**
  - [문서번호] (회사명): 핵심 내용"""
        st.code(notebooklm_prompt, language="text")

    with st.expander("📝 처리 로그", expanded=False):
        st.text(st.session_state.log_text)
